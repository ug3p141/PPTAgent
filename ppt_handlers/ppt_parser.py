import os
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from utils import Config, object_to_dict
from collections import Counter

base_config = Config()


# 这个太难了我感觉，先放着，以后有时间再看
def process_chart_shape(shape, frame_style, page_idx, shape_idx):
    chart = shape.chart
    chart_type = chart.chart_type._member_name
    chart_data = []
    for series in chart.series:
        chart_item_data = {"name": str(series.name)}
        data_list = []
        for point in series.values:
            data_list.append(point)
        chart_item_data["data_list"] = data_list
        chart_data.append(chart_item_data)
    return chart_type, chart_data


# 三种级别的text 可以保存哪些属性
# textframe
# paragraph: space, alignment, level, font
# run: font, hyperlink, text
# 以paragraphs 为单位处理文本框
def process_text_shape(shape, frame_style, is_transition_page, page_idx, shape_idx):
    text_frame = shape.text_frame
    frame_style |= object_to_dict(text_frame)
    frame_data = list()
    # 获取文本框内容
    for _, paragraph in enumerate(text_frame.paragraphs):
        font_counter = Counter()
        text = (
            paragraph.text
        )  # 有些shape的text为空但也有用（例如一些矩形框），所以不对text做判断
        is_transition_page |= True if re.match(r"目\s*录|提\s*纲", text) else False
        content = dict()
        content["type"] = "text"
        content["level"] = paragraph.level
        # 只有fill属性懒得解析了
        content["font"] = Counter([run.font for run in paragraph.runs])
        object_to_dict(paragraph.font) | {"size": paragraph.font.size}
        try:
            content["font"]["color"] = paragraph.font.color.rgb
        except:
            content["font"]["color"] = None

        # element_style = {}
        # 获取Run对象的字体样式 由于runs的信息负责，我认为此处应该使用markdown进行解析，放置过多复杂信息
        # assert len(paragraph.runs) == 1
        # for run in paragraph.runs:
        #     font = run.font
        #     font_style = object_to_dict(font)
        #     try:
        #         font_style['color'] = font.color.rgb
        #     except:
        #         font_style["font_color"] = None
        # content["element_style"] = element_style

        content["text"] = text  # 稍后做markdown处理
        frame_data.append(content)
    return frame_data, is_transition_page


# 以一个markdown解析吧后续
def process_table_shape(shape, frame_style, page_idx, shape_idx):
    frame_style.update(object_to_dict(shape.table))
    table = shape.table
    row_list = list()
    for row in table.rows:
        cell_list = list()
        for cell in row.cells:
            # 获取每个单元格的内容
            cell_data = cell.text_frame.text
            cell_list.append(cell_data)
        row_list.append(cell_list)
    return row_list


def process_image_shape(shape, frame_style, file_path, shape_id):
    fname, _ = os.path.splitext(os.path.basename(file_path))
    img_path = os.path.join(
        base_config.IMAGE_DIR, f"{fname}-{shape_id}.{shape.image.ext}"
    )
    frame_style.update(
        {
            "dpi": shape.image.dpi,
            "filename": shape.image.filename,
            "size": shape.image.size,
            "crop_bottom": shape.crop_bottom,
            "crop_top": shape.crop_top,
            "crop_left": shape.crop_left,
            "crop_right": shape.crop_right,
            "auto_shape_type": shape.auto_shape_type,
        }
    )
    # 保存图像文件
    with open(img_path, "wb") as f:
        f.write(shape.image.blob)
    frame_data = os.path.basename(img_path)
    return frame_data


def parser_ppt(file_path):
    # 打开PPT文件
    ppt = Presentation(file_path)
    file_name = os.path.basename(file_path)

    rs_data = dict()
    rs_data["total_page"] = len(ppt.slides)
    rs_data["author"] = ppt.core_properties.author
    rs_data["file_name"] = file_name
    rs_data["slide_width"] = ppt.slide_width.cm
    rs_data["slide_height"] = ppt.slide_height.cm
    rs_data["slide_pages"] = []
    for page_idx, slide in enumerate(ppt.slides):
        page_data = {
            "slide_idx": page_idx,
            "slide_layout_name": (
                slide.slide_layout.name if slide.slide_layout else None
            ),
            "slide_title": slide.shapes.title.text if slide.shapes.title else None,
            "transition_page": False,
        }
        slide_data = list()
        # 将组合形状展平
        shape_list = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:  # 此处无父子关系只有组合
                for sub_shape in shape.shapes:
                    shape_list.append(sub_shape)
            else:
                shape_list.append(shape)
        # 遍历每个形状（文本框、图像等）
        is_transition_page = False  # 是否是过渡页 目录提纲等
        for shape_idx, shape in enumerate(shape_list):
            frame_dict = dict()
            frame_style = {
                "top": shape.top.cm,
                "left": shape.left.cm,
                "width": shape.width.cm,
                "height": shape.height.cm,
                "rotation": shape.rotation,
            }
            frame_data, frame_type = None, ""

            if shape.has_text_frame:  # 获取文本框内容
                frame_type = "text"
                frame_data, is_transition_page = process_text_shape(
                    shape, frame_style, is_transition_page, page_idx, shape_idx
                )
            elif shape.has_chart:  # 获取框图内容
                frame_type = "chart"
                frame_dict["chart_type"], frame_data = process_chart_shape(
                    shape, frame_style, page_idx, shape_idx
                )
            elif shape.has_table:  # 保存表格数据
                frame_type = "table"
                frame_data = process_table_shape(
                    shape, frame_style, page_idx, shape_idx
                )
            elif "picture" in str(shape.shape_type).lower():  # 图片
                frame_type = "picture"
                frame_data = process_image_shape(
                    shape, frame_style, file_path, shape_idx
                )
            elif shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                print(shape.shape_type, "其他类型shape，请检查", page_idx, shape_idx)

            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:  # 自动形状
                frame_type = "auto_shape"
                frame_dict["shape_type"] = str(shape.auto_shape_type)

            if frame_data:
                frame_dict.update(
                    {
                        "frame_style": frame_style,
                        "frame_data": frame_data,
                        "type": frame_type,
                    }
                )
                slide_data.append(frame_dict)
        slide_data.sort(
            key=lambda k: (
                k.get("frame_style").get("top"),
                k.get("frame_style").get("left"),
            ),
            reverse=False,
        )
        page_data["transition_page"] = is_transition_page
        page_data["slide_data"] = slide_data
        rs_data["slide_pages"].append(page_data)

    return rs_data


if __name__ == "__main__":
    json_data = parser_ppt(
        os.path.join(
            base_config.PPT_DIR,
            "中文信息联合党支部2022年述职报告.pptx",
        )
    )
