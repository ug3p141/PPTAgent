import os
from utils import Config, dict_to_object
from ppt_parser import parser_ppt
from pptx import Presentation
from pptx.util import Cm

from pptx.dml.color import RGBColor

base_config = Config()


def ppt_generate(ppt_data, out_f_name, ppt_template: dict = None):
    ppt = Presentation(ppt_template)
    ppt.slide_width = Cm(ppt_data["slide_width"])
    ppt.slide_height = Cm(ppt_data["slide_height"])

    # 设置 PPT 的元数据
    ppt.core_properties.author = "OminiPreGen"
    slide_layouts = {sl.name: sl for sl in ppt.slide_layouts}
    default_layout = list(slide_layouts.values())[0]

    # 遍历每个幻灯片
    for slide_data in ppt_data["slide_pages"]:
        # 添加一个新的幻灯片
        slide_layout = slide_layouts.get(slide_data["slide_layout_name"], default_layout)
        slide = ppt.slides.add_slide(slide_layout)

        # 遍历每个形状
        for frame in slide_data["slide_data"]:
            if frame["type"] == "text":
                # 添加一个文本框
                left = Cm(frame["frame_style"]["left"])
                top = Cm(frame["frame_style"]["top"])
                width = Cm(frame["frame_style"]["width"])
                height = Cm(frame["frame_style"]["height"])
                txBox = slide.shapes.add_textbox(left, top, width, height)

                # 在文本框中添加文本
                tf = txBox.text_frame
                for para in frame["frame_data"]:
                    p = tf.add_paragraph()
                    p.level = para["level"]
                    p.text = para["text"]
                    font_color = para["font"].pop("color")
                    if font_color is not None:
                        p.font.color.rgb = RGBColor(font_color)
                    dict_to_object(para["font"], p.font)

            elif frame["type"] == "picture":
                # 添加一个图片
                left = Cm(frame["frame_style"]["left"])
                top = Cm(frame["frame_style"]["top"])
                width = Cm(frame["frame_style"]["width"])
                height = Cm(frame["frame_style"]["height"])
                slide.shapes.add_picture(
                    os.path.join(base_config.IMAGE_DIR, frame["frame_data"]),
                    left,
                    top,
                    width,
                    height,
                )

    # 保存 PPT 文件
    ppt.save(os.path.join(base_config.GEN_PPT_DIR, out_f_name))


if __name__ == "__main__":
    ppt_path = os.path.join(
        base_config.PPT_DIR, "中文信息联合党支部2022年述职报告.pptx"
    )
    ppt_data = parser_ppt(ppt_path)

    ppt_generate(ppt_data, "test_ppt_parse&gen.pptx", ppt_path)
