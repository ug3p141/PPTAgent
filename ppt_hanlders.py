from collections import Counter
import os
import re
from rich import print
from pptx.oxml import parse_xml
from pptx import Presentation as PPTXPre
from pptx.slide import Slide
from pptx.shapes.autoshape import Shape as PPTXAutoShape
from pptx.shapes.picture import Picture
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.chart.chart import Chart
from pptx.table import Table
from pptx.text.text import _Run, TextFrame
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing_extensions import Dict, Type
from utils import Config, dict_to_object, xml_print, object_to_dict  # noqa
from pptx.dml.color import RGBColor

base_config = Config()
pjoin = os.path.join


# 三种级别的text 可以保存的属性
# textframe: shape bounds
# paragraph: space, alignment, level, font
# run: font, hyperlink, text
# 以paragraphs 为单位处理文本框
class TextFrame:

    def __init__(
        self,
        is_textframe: bool,
        style: Dict,
        data: Dict,
    ):
        self.is_textframe = is_textframe
        self.style = style
        self.data = data

    @classmethod
    def from_shape(cls, shape: BaseShape):
        if not shape.has_text_frame:
            return cls(False, {}, [])
        shape = shape.text_frame
        style = object_to_dict(shape)
        style.pop("text")
        data = []
        # 获取文本框内容
        for _, paragraph in enumerate(shape.paragraphs):
            # 有些shape的text为空但也有用（例如一些矩形框），所以不对text做判断
            if "支部基本情况(党员队伍建设)" in paragraph.text:
                pass
            if not paragraph.text:
                continue
            runs = paragraph.runs
            if len(runs) == 0:  # 成功解析
                runs = [
                    _Run(r, paragraph)
                    for r in parse_xml(paragraph._element.xml.replace("fld", "r")).r_lst
                ]
            # 只有fill属性懒得解析了
            content = {
                "font": object_to_dict(
                    Counter([run.font for run in runs]).most_common(1)[0][0]
                ),
            } | object_to_dict(paragraph)
            content.pop("runs", None)
            if content["font"]["name"] == "+mj-ea":
                content["font"]["name"] = "宋体"

            data.append(content)
            # element_style = {}
            # 获取Run对象的字体样式 由于runs的信息负责，我认为此处应该使用markdown进行解析，放置过多复杂信息
            # assert len(paragraph.runs) == 1
            # for run in paragraph.runs:
            #     font = run.font
            #     font_style = object_to_dict(font)
            #     try:
            #         font_style['color'] = font.color.rgb
            #     except:
            #         font_style["font_color"] = None
            # content["element_style"] = element_style
        return cls(True, style, data)

    def build(self, shape: BaseShape):
        if not self.is_textframe:
            return
        tf = (
            shape.text_frame
            if shape.has_text_frame
            else shape.add_textbox(**self.style["shape_bounds"])
        )
        # 注意，只对textframe是无效的，还要设置一遍所有的runs
        for para in self.data:
            p = tf.add_paragraph()
            run = p.add_run()
            font = para.pop("font")
            if font["color"] is not None:
                run.font.fill.solid()
                run.font.fill.fore_color.rgb = RGBColor(*font.pop("color"))
            dict_to_object(para, run)
            dict_to_object(font, run.font)


# TODO 判断build后的属性是否与解析前的一直


class UnsupportedShape:
    @classmethod
    def from_shape(
        cls, shape_idx: int, shape: BaseShape, style: Dict, text_frame: TextFrame
    ):
        print(f"Unsupport Shape --- {shape.__class__}\n", object_to_dict(shape))
        assert not any([shape.has_chart, shape.has_table, shape.has_text_frame])
        return cls()

    def build(self, slide: Slide):
        pass


class ShapeElement:
    def __init__(
        self,
        shape_idx: int,
        style: Dict,
        data: Dict,
        text_frame: TextFrame,
    ):
        self.shape_idx = shape_idx
        self.style = style
        self.data = data
        self.text_frame = text_frame
        self.validation = True

    @classmethod
    def from_shape(cls, shape_idx: int, shape: BaseShape):
        style = {
            "shape_bounds": {
                "width": shape.width,
                "height": shape.height,
                "left": shape.left,
                "top": shape.top,
            },
            "shape_type": shape.shape_type,
            "rotation": shape.rotation,
        }
        text_frame = TextFrame.from_shape(shape)
        return SHAPECAST.get(shape.shape_type, UnsupportedShape).from_shape(
            shape_idx, shape, style, text_frame
        )

    def get_style(self):
        return self.style

    def get_data(self):
        return self.data

    def set_validation(self, validation: bool):
        self.validation = validation

    def build(self, slide, shape):
        self.text_frame.build(shape)
        if "rotation" in self.style:
            shape.rotation = self.style["rotation"]
        if self.validation:
            new_shape = ShapeElement.from_shape(self.shape_idx, shape)
            assert self.data == new_shape.data
            assert self.style == new_shape.style


class AutoShape(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        shape_idx: int,
        shape: PPTXAutoShape,
        style: Dict,
        text_frame: TextFrame,
    ):
        data = {
            "auto_shape_type": shape.auto_shape_type,
            # "fill_color": shape.fill.fore_color.rgb,
        }
        return cls(shape_idx, style, data, text_frame)

    def build(self, slide: Slide):
        shape = slide.shapes.add_shape(
            self.data["auto_shape_type"], **self.style["shape_bounds"]
        )
        super().build(slide, shape)


class TextBox(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        shape_idx: int,
        shape: TextFrame,
        style: Dict,
        text_frame: TextFrame,
    ):
        return cls(shape_idx, style, None, text_frame)

    def build(self, slide: Slide):
        shape = slide.shapes.add_textbox(**self.style["shape_bounds"])
        super().build(slide, shape)


class Placeholder(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        shape_idx: int,
        shape: SlidePlaceholder,
        style: Dict,
        text_frame: TextFrame,
    ):

        style |= {
            "placeholder_type": shape.placeholder_format.type,
        }
        if not shape.has_text_frame or shape.has_chart or shape.has_table:
            print(f"Unsupported Shape: {shape}")
        data = []
        return cls(shape_idx, style, data, text_frame)

    def build(self, slide: Slide):
        pass
        # super().build(slide, shape)


class Picture(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        shape_idx: int,
        shape: Picture,
        style: Dict,
        text_frame: TextFrame,
    ):
        img_name = re.sub(r"[\/\0]", "_", shape.name)
        img_path = pjoin(
            base_config.IMAGE_DIR,
            f"{img_name}.{shape.image.ext}",
        )
        style["img_style"] = {
            "crop_bottom": shape.crop_bottom,
            "crop_top": shape.crop_top,
            "crop_left": shape.crop_left,
            "crop_right": shape.crop_right,
            "auto_shape_type": shape.auto_shape_type,
        }
        with open(img_path, "wb") as f:
            f.write(shape.image.blob)
        return cls(shape_idx, style, [img_path, shape.name], text_frame)

    def build(self, slide: Slide):
        shape = slide.shapes.add_picture(
            self.data[0],
            **self.style["shape_bounds"],
        )
        shape.name = self.data[1]
        dict_to_object(self.style["img_style"], shape.image)
        super().build(slide, shape)


class GroupShape(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        shape_idx: int,
        shape: GroupShape,
        style: Dict,
        text_frame: TextFrame,
    ):
        data = [
            ShapeElement.from_shape(i, sub_shape)
            for i, sub_shape in enumerate(shape.shapes)
        ]
        return cls(shape_idx, style, data, text_frame)

    def build(self, slide: Slide):
        for sub_shape in self.data:
            sub_shape.build(slide)
        # super().build(slide, shape)


class Chart(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        shape_idx: int,
        shape: Chart,
        style: Dict,
        text_frame: TextFrame,
    ):
        chart = shape.chart
        style["chart_style"] = object_to_dict(chart) | {
            "chart_title": (
                chart.chart_title.text_frame.text
                if chart.chart_title.has_text_frame
                else None
            ),
        }
        chart_data = []
        for series in chart.series:
            chart_item_data = {"name": str(series.name)}
            data_list = []
            for point in series.values:
                data_list.append(point)
            chart_item_data["data_list"] = data_list
            chart_data.append(chart_item_data)
        obj = cls(shape_idx, style, chart_data, text_frame)
        setattr(obj, "shape", shape)
        return obj

    def build(self, slide: Slide):
        pass
        # graphic_frame = slide.shapes.add_chart(
        #     self.style["chart_style"]["chart_type"],

        # shape
        # super().build(slide, shape)


# 目前的问题，rows和cols的宽度或高度获取不到
class Table(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        shape_idx: int,
        shape: Table,
        style: Dict,
        text_frame: TextFrame,
    ):
        table = shape.table
        row_list = list()
        for row in table.rows:
            cell_list = list()
            for cell in row.cells:
                cell_data = cell.text_frame.text
                cell_list.append(cell_data)
            row_list.append(cell_list)
        style["table_style"] = object_to_dict(shape.table)

        obj = cls(shape_idx, style, row_list, text_frame)
        setattr(obj, "shape", shape)
        return obj

    def build(self, slide: Slide):
        graphic_frame = slide.shapes.add_table(
            len(self.data), len(self.data[0]), **self.style["shape_bounds"]
        )

        dict_to_object(self.style["table_style"], graphic_frame.table)
        self.text_frame.build(graphic_frame)


# template包含哪些东西，如何定义一个template
# 1. template pics：例如banner、background pic、logo，此部分只需要手动识别插入以及build
# 2. cloneable shapes：例如listed text, parallel text等，此部分需要手动安排布局
class SlidePage:
    def __init__(
        self,
        shapes: list[ShapeElement],
        slide_idx: int,
        slide_notes: str,
        slide_layout_name: str,
        slide_title: str,
    ):
        self.slide_idx = slide_idx
        self.shapes = shapes
        self.slide_layout_name = slide_layout_name
        self.slide_title = slide_title

    @classmethod
    def from_slide(cls, slide: Slide, slide_idx: int):
        shapes = [
            ShapeElement.from_shape(i, shape) for i, shape in enumerate(slide.shapes)
        ]
        slide_layout_name = slide.slide_layout.name if slide.slide_layout else None
        # is shape a title
        slide_title = slide.shapes.title.text if slide.shapes.title else None
        slide_notes = (
            slide.notes_slide.notes_text_frame.text
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame
            else None
        )
        return cls(shapes, slide_idx, slide_notes, slide_layout_name, slide_title)

    def build(self, prs: PPTXPre, slide_layout):
        slide = prs.slides.add_slide(slide_layout)
        for shape in self.shapes:
            shape.build(slide)
        # if self.slide_title:
        # slide.shapes.title.text = self.slide_title


class Presentation:
    def __init__(
        self,
        slides: list[SlidePage],
        slide_width: float,
        slide_height: float,
        num_pages: int,
        author: str,
    ) -> None:
        self.slides = slides
        self.slide_width = slide_width
        self.slide_height = slide_height
        self.num_pages = num_pages
        self.author = author
        self.consumed = False

    @classmethod
    def from_file(cls, file_path):
        prs = PPTXPre(file_path)
        slides = [SlidePage.from_slide(slide, i) for i, slide in enumerate(prs.slides)]
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        core_properties = prs.core_properties
        num_pages = len(slides)
        author = core_properties.author
        return cls(slides, slide_width, slide_height, num_pages, author)

    def save(self, file_path):
        if self.consumed:
            raise Exception("Presentation has been consumed")
        self.consumed = True
        prs = PPTXPre()
        prs.core_properties.author = "OminiPreGen" + self.author
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        blank_slide_layout = prs.slide_layouts[6]
        for slide in self.slides:
            slide.build(prs, blank_slide_layout)
        prs.save(file_path)


SHAPECAST: Dict[int, Type[ShapeElement]] = {
    MSO_SHAPE_TYPE.AUTO_SHAPE: AutoShape,
    MSO_SHAPE_TYPE.PLACEHOLDER: Placeholder,
    MSO_SHAPE_TYPE.PICTURE: Picture,
    MSO_SHAPE_TYPE.GROUP: GroupShape,
    MSO_SHAPE_TYPE.CHART: Chart,
    MSO_SHAPE_TYPE.TABLE: Table,
    MSO_SHAPE_TYPE.TEXT_BOX: TextBox,
}


if __name__ == "__main__":
    Presentation.from_file(
        pjoin(
            base_config.PPT_DIR,
            "中文信息联合党支部2022年述职报告.pptx",
        )
    ).save(pjoin(base_config.GEN_PPT_DIR, "ppt_handlers_test.pptx"))
