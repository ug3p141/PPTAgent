import asyncio
import os
import shutil
import subprocess
import tempfile
import traceback
from itertools import product
from time import sleep, time
from types import SimpleNamespace
from typing import Dict, List, Optional, Set, Any
import logging

from html2image import Html2Image
import json_repair
import Levenshtein
from mistune import html as markdown
from PIL import Image as PILImage
from shutil import which
from pdf2image import convert_from_path
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.shapes.group import GroupShape
from pptx.text.text import _Paragraph, _Run
from pptx.util import Length, Pt
from pptx.parts.image import Image
from tenacity import RetryCallState, retry, stop_after_attempt, wait_fixed


def get_logger(name="pptagent", level=None):
    """
    Get a logger with the specified name and level.

    Args:
        name (str): The name of the logger.
        level (int): The logging level (default: logging.INFO).

    Returns:
        logging.Logger: A configured logger instance.
    """
    if level is None:
        level = os.environ.get("LOG_LEVEL", logging.INFO)

    logger = logging.getLogger(name)
    logger.setLevel(level)

    # Check if the logger already has handlers to avoid duplicates
    if not logger.handlers:
        # Create console handler and set level
        console_handler = logging.StreamHandler()
        console_handler.setLevel(level)

        # Create formatter
        formatter = logging.Formatter(
            "%(asctime)s - %(levelname)s - %(filename)s:%(lineno)d - %(message)s"
        )
        console_handler.setFormatter(formatter)

        # Add handler to logger
        logger.addHandler(console_handler)

    return logger


logger = get_logger(__name__)

if which("soffice") is None:
    logging.warning("soffice is not installed, pptx to images conversion will not work")

# Set of supported image extensions
IMAGE_EXTENSIONS: Set[str] = {
    "bmp",
    "jpg",
    "jpeg",
    "pgm",
    "png",
    "ppm",
    "tif",
    "tiff",
    "webp",
}

# Common colors and measurements
BLACK = RGBColor(0, 0, 0)
YELLOW = RGBColor(255, 255, 0)
BLUE = RGBColor(0, 0, 255)
BORDER_LEN = Pt(2)
BORDER_OFFSET = Pt(2)
LABEL_LEN = Pt(24)
FONT_LEN = Pt(20)


def is_image_path(file: str) -> bool:
    """
    Check if a file path is an image based on its extension.

    Args:
        file (str): The file path to check.

    Returns:
        bool: True if the file is an image, False otherwise.
    """
    return file.split(".")[-1].lower() in IMAGE_EXTENSIONS


def get_font_style(font: Dict[str, Any]) -> str:
    """
    Convert a font dictionary to a CSS style string.

    Args:
        font (Dict[str, Any]): The font dictionary.

    Returns:
        str: The CSS style string.
    """
    font = SimpleNamespace(**font)
    styles = []

    if hasattr(font, "size") and font.size:
        styles.append(f"font-size: {font.size}pt")

    if hasattr(font, "color") and font.color:
        if all(c in "0123456789abcdefABCDEF" for c in font.color):
            styles.append(f"color: #{font.color}")
        else:
            styles.append(f"color: {font.color}")

    if hasattr(font, "bold") and font.bold:
        styles.append("font-weight: bold")

    if hasattr(font, "italic") and font.italic:
        styles.append("font-style: italic")

    return "; ".join(styles)


def runs_merge(paragraph: _Paragraph) -> Optional[_Run]:
    """
    Merge all runs in a paragraph into a single run.

    Args:
        paragraph (_Paragraph): The paragraph to merge runs in.

    Returns:
        Optional[_Run]: The merged run, or None if there are no runs.
    """
    runs = paragraph.runs

    # Handle field codes
    if len(runs) == 0:
        runs = [
            _Run(r, paragraph)
            for r in parse_xml(paragraph._element.xml.replace("fld", "r")).r_lst
        ]
    if len(runs) == 1:
        return runs[0]
    if len(runs) == 0:
        return None

    # Find the run with the most text
    run = max(runs, key=lambda x: len(x.text))
    run.text = paragraph.text

    # Remove other runs
    for r in runs:
        if r != run:
            r._r.getparent().remove(r._r)
    return run


def older_than(filepath: str, seconds: int = 10, wait: bool = False) -> bool:
    """
    Check if a file is older than a specified number of seconds.

    Args:
        filepath (str): The path to the file.
        seconds (int): The number of seconds to check against.
        wait (bool): Whether to wait for the file to exist.

    Returns:
        bool: True if the file is older than the specified number of seconds, False otherwise.
    """
    if not os.path.exists(filepath):
        while wait:
            logger.info("waiting for: %s", filepath)
            sleep(1)
            if os.path.exists(filepath):
                sleep(seconds)
                return True
        return False
    file_creation_time = os.path.getctime(filepath)
    current_time = time()
    return seconds < (current_time - file_creation_time)


def edit_distance(text1: str, text2: str) -> float:
    """
    Calculate the normalized edit distance between two strings.

    Args:
        text1 (str): The first string.
        text2 (str): The second string.

    Returns:
        float: The normalized edit distance (0.0 to 1.0, where 1.0 means identical).
    """
    if not text1 and not text2:
        return 1.0
    return 1 - Levenshtein.distance(text1, text2) / max(len(text1), len(text2))


def tenacity_log(retry_state: RetryCallState) -> None:
    """
    Log function for tenacity retries.

    Args:
        retry_state (RetryCallState): The retry state.
    """
    logger.warning("tenacity retry: %s", retry_state)
    traceback.print_tb(retry_state.outcome.exception().__traceback__)


def get_json_from_response(response: str) -> Dict[str, Any]:
    """
    Extract JSON from a text response.

    Args:
        response (str): The response text.

    Returns:
        Dict[str, Any]: The extracted JSON.

    Raises:
        Exception: If JSON cannot be extracted from the response.
    """
    response = response.strip()

    # Try to extract JSON from markdown code blocks
    l, r = response.rfind("```json"), response.rfind("```")
    if l != -1 and r != -1:
        json_obj = json_repair.loads(response[l + 7 : r].strip())
        if isinstance(json_obj, (dict, list)):
            return json_obj

    # Try to find JSON by looking for matching braces
    open_braces = []
    close_braces = []

    for i, char in enumerate(response):
        if char == "{":
            open_braces.append(i)
        elif char == "}":
            close_braces.append(i)

    for i, j in product(open_braces, reversed(close_braces)):
        if i > j:
            continue
        try:
            json_obj = json_repair.loads(response[i : j + 1])
            if isinstance(json_obj, (dict, list)):
                return json_obj
        except Exception:
            pass

    raise Exception("JSON not found in the given output", response)


# Create a tenacity decorator with custom settings
tenacity = retry(
    wait=wait_fixed(3), stop=stop_after_attempt(5), after=tenacity_log, reraise=True
)


def split_markdown_by_level(text: str, level: int = 1):
    """
    Split a markdown document into sections based on the level of headings.

    Args:
        text (str): The markdown document to split.
        level (int): The level of headings to split the document into.

    Returns:
        List[Tuple[str, str]]: A list of tuples, each containing a header and content.
    """
    lines = text.splitlines()
    current_header = None
    current_content = []

    for line in lines:
        if line.strip().startswith("#" * level + " "):
            if sum(len(line) for line in current_content) > 128:
                yield current_header, "\n".join(current_content)
                current_header = line
                current_content = []
            else:
                current_content.append(line)
        else:
            current_content.append(line)

    if current_content:
        yield current_header, "\n".join(current_content)


def split_markdown_to_chunks(
    markdown_text: str, max_length: int = 16384, max_level: int = 3
) -> List[Dict[str, str]]:
    """
    Split a markdown document into chunks of a maximum length.

    Args:
        markdown_text (str): The markdown document to split.
        max_length (int): The maximum length of a chunk.
        max_level (int): The maximum level of headings to split the document into.

    Returns:
        List[Dict[str, str]]: A list of dictionaries, each containing a header and content.

    Raises:
        ValueError: If the markdown document is too long to be split into chunks.
    """
    for level in range(1, max_level + 1):
        sections = []
        for header, content in split_markdown_by_level(markdown_text, level):
            if len(content) > max_length:
                sections = []
                break
            sections.append({"header": header, "content": content})
        if len(sections) != 0:
            return sections
    raise ValueError(
        f"The markdown document is too long to be split into chunks, with max_level = {max_level} and max_length = {max_length}"
    )


TABLE_CSS = """
table {
    border-collapse: collapse;  /* Merge borders */
    width: auto;               /* Width adapts to content */
    font-family: SimHei, Arial, sans-serif;  /* Font supporting Chinese characters */
    background: white;
}
th, td {
    border: 1px solid black;  /* Add borders */
    padding: 8px;             /* Cell padding */
    text-align: center;       /* Center text */
}
th {
    background-color: #f2f2f2; /* Header background color */
}
"""


# Convert Markdown to HTML
def markdown_table_to_image(markdown_text: str, output_path: str):
    """
    Convert a Markdown table to a cropped image

    Args:
    markdown_text (str): Markdown text containing a table
    output_path (str): Output image path, defaults to 'table_cropped.png'

    Returns:
    str: The path of the generated image
    """
    html = markdown(markdown_text, extensions=["tables"])
    assert "table" in html, "Failed to find table in markdown"

    parent_dir, basename = os.path.split(output_path)
    hti = Html2Image(disable_logging=True, output_path=parent_dir)
    hti.browser.use_new_headless = None
    hti.screenshot(html_str=html, css_str=TABLE_CSS, save_as=basename)

    img = PILImage.open(output_path).convert("RGB")
    bbox = img.getbbox()
    assert (
        bbox is not None
    ), "Failed to capture the bbox, may be markdown table conversion failed"
    bbox = (0, 0, bbox[2] + 10, bbox[3] + 10)
    img.crop(bbox).save(output_path)
    return output_path


@tenacity
def ppt_to_images(file: str, output_dir: str):
    assert pexists(file), f"File {file} does not exist"
    if pexists(output_dir):
        logger.warning(f"ppt2images: {output_dir} already exists")
    os.makedirs(output_dir, exist_ok=True)
    with tempfile.TemporaryDirectory() as temp_dir:
        command_list = [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            file,
            "--outdir",
            temp_dir,
        ]
        subprocess.run(command_list, check=True, stdout=subprocess.DEVNULL)

        for f in os.listdir(temp_dir):
            if not f.endswith(".pdf"):
                continue
            temp_pdf = pjoin(temp_dir, f)
            images = convert_from_path(temp_pdf, dpi=72)
            for i, img in enumerate(images):
                img.save(pjoin(output_dir, f"slide_{i+1:04d}.jpg"))
            return

        raise RuntimeError("No PDF file was created in the temporary directory", file)


async def ppt_to_images_async(file: str, output_dir: str):
    assert pexists(file), f"File {file} does not exist"
    if pexists(output_dir):
        logger.warning(f"ppt2images: {output_dir} already exists")
    os.makedirs(output_dir, exist_ok=True)

    with tempfile.TemporaryDirectory() as temp_dir:
        command_list = [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            file,
            "--outdir",
            temp_dir,
        ]

        process = await asyncio.create_subprocess_exec(
            *command_list,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        _, stderr = await process.communicate()
        if process.returncode != 0:
            raise RuntimeError(f"soffice failed with error: {stderr.decode()}")
        for f in os.listdir(temp_dir):
            if not f.endswith(".pdf"):
                continue
            temp_pdf = pjoin(temp_dir, f)
            images = convert_from_path(temp_pdf, dpi=72)
            for i, img in enumerate(images):
                img.save(pjoin(output_dir, f"slide_{i+1:04d}.jpg"))
            return

        raise RuntimeError(
            f"No PDF file was created in the temporary directory: {file}"
        )


def parsing_image(image: Image, image_path: str) -> str:
    # Handle WMF images (PDFs)
    if image.ext == "wmf":
        image_path = image_path.replace(".wmf", ".jpg")
        if not pexists(image_path):
            wmf_to_images(image.blob, image_path)
    # Check for supported image types
    elif image.ext not in IMAGE_EXTENSIONS:
        raise ValueError(f"Unsupported image type {image.ext}")

    # Save image if it doesn't exist
    if not pexists(image_path):
        with open(image_path, "wb") as f:
            f.write(image.blob)
    return image_path


@tenacity
def wmf_to_images(blob: bytes, filepath: str):
    if not filepath.endswith(".jpg"):
        raise ValueError("filepath must end with .jpg")
    dirname = os.path.dirname(filepath)
    basename = os.path.basename(filepath).removesuffix(".jpg")
    with tempfile.TemporaryDirectory() as temp_dir:
        with open(pjoin(temp_dir, f"{basename}.wmf"), "wb") as f:
            f.write(blob)
        command_list = [
            "soffice",
            "--headless",
            "--convert-to",
            "jpg",
            pjoin(temp_dir, f"{basename}.wmf"),
            "--outdir",
            dirname,
        ]
        subprocess.run(command_list, check=True, stdout=subprocess.DEVNULL)

    assert pexists(filepath), f"File {filepath} does not exist"


def parse_groupshape(groupshape: GroupShape) -> List[Dict[str, Length]]:
    """
    Parse a group shape to get the bounds of its child shapes.

    Args:
        groupshape (GroupShape): The group shape to parse.

    Returns:
        List[Dict[str, Length]]: The bounds of the child shapes.

    Raises:
        AssertionError: If the input is not a GroupShape.
    """
    assert isinstance(groupshape, GroupShape), "Input must be a GroupShape"

    # Get group bounds
    group_top_left_x = groupshape.left
    group_top_left_y = groupshape.top
    group_width = groupshape.width
    group_height = groupshape.height

    # Get shape bounds
    shape_top_left_x = min([sp.left for sp in groupshape.shapes])
    shape_top_left_y = min([sp.top for sp in groupshape.shapes])
    shape_width = (
        max([sp.left + sp.width for sp in groupshape.shapes]) - shape_top_left_x
    )
    shape_height = (
        max([sp.top + sp.height for sp in groupshape.shapes]) - shape_top_left_y
    )

    # Calculate bounds for each shape in the group
    group_shape_xy = []
    for sp in groupshape.shapes:
        group_shape_left = (
            sp.left - shape_top_left_x
        ) * group_width / shape_width + group_top_left_x
        group_shape_top = (
            sp.top - shape_top_left_y
        ) * group_height / shape_height + group_top_left_y
        group_shape_width = sp.width * group_width / shape_width
        group_shape_height = sp.height * group_height / shape_height

        group_shape_xy.append(
            {
                "left": Length(group_shape_left),
                "top": Length(group_shape_top),
                "width": Length(group_shape_width),
                "height": Length(group_shape_height),
            }
        )

    return group_shape_xy


def is_primitive(obj: Any) -> bool:
    """
    Check if an object is a primitive type or a collection of primitive types.

    Args:
        obj (Any): The object to check.

    Returns:
        bool: True if the object is a primitive type or a collection of primitive types, False otherwise.
    """
    if isinstance(obj, (list, tuple, set, frozenset)):
        return all(is_primitive(item) for item in obj)

    return isinstance(
        obj, (int, float, complex, bool, str, bytes, bytearray, type(None))
    )


DEFAULT_EXCLUDE: Set[str] = set(["element", "language_id", "ln", "placeholder_format"])


def object_to_dict(
    obj: Any,
    result: Optional[Dict[str, Any]] = None,
    exclude: Optional[Set[str]] = None,
) -> Dict[str, Any]:
    """
    Convert an object to a dictionary.

    Args:
        obj (Any): The object to convert.
        result (Optional[Dict[str, Any]]): The result dictionary to update.
        exclude (Optional[Set[str]]): The attributes to exclude.

    Returns:
        Dict[str, Any]: The dictionary representation of the object.
    """
    if result is None:
        result = {}

    exclude = DEFAULT_EXCLUDE.union(exclude or set())

    for attr in dir(obj):
        if attr in exclude or attr.startswith("_") or callable(getattr(obj, attr)):
            continue

        try:
            attr_value = getattr(obj, attr)

            # Handle complex numbers
            if hasattr(attr_value, "real"):
                attr_value = attr_value.real

            # Convert size to points
            if attr == "size" and isinstance(attr_value, int):
                attr_value = Length(attr_value).pt

            # Only include primitive types
            if is_primitive(attr_value):
                result[attr] = attr_value
        except Exception:
            pass

    return result


def merge_dict(d1: Dict[str, Any], d2: List[Dict[str, Any]]) -> Dict[str, Any]:
    """
    Merge multiple dictionaries, keeping only values that are the same across all dictionaries.

    Args:
        d1 (Dict[str, Any]): The base dictionary.
        d2 (List[Dict[str, Any]]): The dictionaries to merge.

    Returns:
        Dict[str, Any]: The merged dictionary.
    """
    if not d2:
        return d1

    for key in list(d1.keys()):
        values = [d[key] for d in d2]
        if d1[key] is not None and len(values) != 1:
            values.append(d1[key])

        # Skip if any value is None or values are not all the same
        if (
            not values
            or values[0] is None
            or not all(value == values[0] for value in values)
        ):
            continue

        # Set the value in the base dictionary and clear it in the other dictionaries
        d1[key] = values[0]
        for d in d2:
            d[key] = None
    return d1


def dict_to_object(
    dict_obj: Dict[str, Any], obj: Any, exclude: Optional[Set[str]] = None
) -> None:
    """
    Apply dictionary values to an object.

    Args:
        dict_obj (Dict[str, Any]): The dictionary with values to apply.
        obj (Any): The object to apply values to.
        exclude (Optional[Set[str]]): The keys to exclude.
    """
    if exclude is None:
        exclude = set()

    for key, value in dict_obj.items():
        if key not in exclude and value is not None:
            try:
                setattr(obj, key, value)
            except Exception as e:
                logger.error("Error setting attribute %s: %s", key, e)


def package_join(*paths: str) -> str:
    """
    Join paths with the appropriate separator for the platform.

    Args:
        *paths: The paths to join.

    Returns:
        str: The joined path.
    """
    _dir = pdirname(pdirname(__file__))
    return pjoin(_dir, *paths)


class Config:
    """
    Configuration class for the application.
    """

    def __init__(
        self,
        rundir: Optional[str] = None,
        session_id: Optional[str] = None,
        debug: bool = True,
    ):
        """
        Initialize the configuration.

        Args:
            rundir (Optional[str]): The run directory.
            session_id (Optional[str]): The session ID.
            debug (bool): Whether to enable debug mode.
        """
        self.DEBUG = debug

        if rundir is not None:
            self.set_rundir(rundir)
        elif session_id is not None:
            self.set_session(session_id)
        else:
            raise ValueError("No session ID or run directory provided")

    def set_session(self, session_id: str) -> None:
        """
        Set the session ID and update the run directory.

        Args:
            session_id (str): The session ID.
        """
        self.session_id = session_id
        self.set_rundir(f"./runs/{session_id}")

    def set_rundir(self, rundir: str) -> None:
        """
        Set the run directory and create necessary subdirectories.

        Args:
            rundir (str): The run directory.
        """
        self.RUN_DIR = rundir
        self.IMAGE_DIR = pjoin(self.RUN_DIR, "images")

        for the_dir in [self.RUN_DIR, self.IMAGE_DIR]:
            os.makedirs(the_dir, exist_ok=True)

    def set_debug(self, debug: bool) -> None:
        """
        Set the debug mode.

        Args:
            debug (bool): Whether to enable debug mode.
        """
        self.DEBUG = debug

    def remove_rundir(self) -> None:
        """
        Remove the run directory and its subdirectories.
        """
        if pexists(self.RUN_DIR):
            shutil.rmtree(self.RUN_DIR)
        if pexists(self.IMAGE_DIR):
            shutil.rmtree(self.IMAGE_DIR)

    def __repr__(self) -> str:
        """
        Get a string representation of the configuration.

        Returns:
            str: A string representation of the configuration.
        """
        attrs = []
        for attr in dir(self):
            if not attr.startswith("_") and not callable(getattr(self, attr)):
                attrs.append(f"{attr}={getattr(self, attr)}")
        return f"Config({', '.join(attrs)})"


# Path utility functions
pjoin = os.path.join
pexists = os.path.exists
pbasename = os.path.basename
pdirname = os.path.dirname
