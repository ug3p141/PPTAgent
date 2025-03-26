import traceback
from packaging.version import Version
from typing import List, Optional, Tuple, Type, Generator

from pptx import Presentation as PPTXPre
from pptx import __version__ as PPTXVersion
from pptx.slide import Slide as PPTXSlide
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape as PPTXGroupShape

from pptagent.utils import Config, get_logger
from pptagent.shapes import (
    T,
    GroupShape,
    Picture,
    ShapeElement,
    Background,
    StyleArg,
)

PPTXVersion, Mark = PPTXVersion.split("+")
assert (
    Version(PPTXVersion) >= Version("1.0.3") and Mark == "PPTAgent"
), "You should install the version of `python-pptx` maintained specifically for this project."

# Type variable for ShapeElement subclasses

logger = get_logger(__name__)
class SlidePage:
    """
    A class to represent a slide page in a presentation.
    """

    def __init__(
        self,
        shapes: List[ShapeElement],
        backgrounds: List[Background],
        slide_idx: int,
        real_idx: int,
        slide_notes: Optional[str],
        slide_layout_name: Optional[str],
        slide_title: Optional[str],
        slide_width: int,
        slide_height: int,
    ):
        """
        Initialize a SlidePage.

        Args:
            shapes (List[ShapeElement]): The shapes in the slide.
            backgrounds (List[Background]): The backgrounds of the slide.
            slide_idx (int): The index of the slide.
            real_idx (int): The real index of the slide.
            slide_notes (Optional[str]): The notes of the slide.
            slide_layout_name (Optional[str]): The layout name of the slide.
            slide_title (Optional[str]): The title of the slide.
            slide_width (int): The width of the slide.
            slide_height (int): The height of the slide.
        """
        self.shapes = shapes
        self.backgrounds = backgrounds
        self.slide_idx = slide_idx
        self.real_idx = real_idx
        self.slide_notes = slide_notes
        self.slide_layout_name = slide_layout_name
        self.slide_title = slide_title
        self.slide_width = slide_width
        self.slide_height = slide_height

        # Assign group labels to group shapes
        groups_shapes_labels = []
        for shape in self.shape_filter(GroupShape):
            for group_shape in groups_shapes_labels:
                if group_shape == shape:
                    shape.group_label = group_shape.group_label
                    continue
            groups_shapes_labels.append(shape)
            shape.group_label = f"group_{len(groups_shapes_labels)}"

    @classmethod
    def from_slide(
        cls,
        slide: PPTXSlide,
        slide_idx: int,
        real_idx: int,
        slide_width: int,
        slide_height: int,
        config: Config,
    ) -> "SlidePage":
        """
        Create a SlidePage from a PPTXSlide.

        Args:
            slide (PPTXSlide): The slide object.
            slide_idx (int): The index of the slide.
            real_idx (int): The real index of the slide.
            slide_width (int): The width of the slide.
            slide_height (int): The height of the slide.
            config (Config): The configuration object.

        Returns:
            SlidePage: The created SlidePage.
        """
        backgrounds = [Background.from_slide(slide, config)]
        shapes = [
            ShapeElement.from_shape(
                slide_idx, i, shape, config, slide_width * slide_height
            )
            for i, shape in enumerate(slide.shapes)
            if shape.visible
        ]
        for i, s in enumerate(shapes):
            if isinstance(s, Picture) and s.area / s.slide_area > 0.95:
                backgrounds.append(shapes.pop(i))

        slide_layout_name = slide.slide_layout.name if slide.slide_layout else None
        slide_title = slide.shapes.title.text if slide.shapes.title else None
        slide_notes = (
            slide.notes_slide.notes_text_frame.text
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame
            else None
        )

        return cls(
            shapes,
            backgrounds,
            slide_idx,
            real_idx,
            slide_notes,
            slide_layout_name,
            slide_title,
            slide_width,
            slide_height,
        )

    def build(self, slide: PPTXSlide) -> PPTXSlide:
        """
        Build the slide page in a slide.

        Args:
            slide (PPTXSlide): The slide to build the slide page in.

        Returns:
            PPTXSlide: The built slide.
        """
        # Remove existing placeholders
        for ph in slide.placeholders:
            ph.element.getparent().remove(ph.element)

        # Build background
        for background in self.backgrounds:
            background.build(slide)

        # Build shapes and apply closures
        for shape in self.shapes:
            build_shape = shape.build(slide)
            for closure in shape.closures:
                try:
                    closure.apply(build_shape)
                except Exception as e:
                    raise ValueError(f"Failed to apply closures to slides: {e}")
        return slide

    def shape_filter(
        self, shape_type: Type[T], shapes: Optional[List[ShapeElement]] = None
    ) -> Generator[T, None, None]:
        """
        Filter shapes in the slide by type.

        Args:
            shape_type (Type[T]): The type of shapes to filter.
            shapes (Optional[List[ShapeElement]]): The shapes to filter.

        Yields:
            T: The filtered shapes.
        """
        if shapes is None:
            shapes = self.shapes
        for shape in shapes:
            if isinstance(shape, shape_type):
                yield shape
            elif isinstance(shape, GroupShape):
                yield from self.shape_filter(shape_type, shape.data)

    def get_content_type(self) -> str:
        """
        Get the content type of the slide.

        Returns:
            str: The content type of the slide.
        """
        if len(list(self.shape_filter(Picture))) == 0:
            return "text"
        return "image"

    def to_html(self, style_args: Optional[StyleArg] = None, **kwargs) -> str:
        """
        Represent the slide page in HTML.

        Args:
            style_args (Optional[StyleArg]): The style arguments for HTML conversion.
            **kwargs: Additional arguments.

        Returns:
            str: The HTML representation of the slide page.
        """
        if style_args is None:
            style_args = StyleArg(**kwargs)
        return "".join(
            [
                "<!DOCTYPE html>\n<html>\n",
                (f"<title>{self.slide_title}</title>\n" if self.slide_title else ""),
                f'<body style="width:{self.slide_width}pt; height:{self.slide_height}pt;">\n',
                "\n".join([shape.to_html(style_args) for shape in self.shapes]),
                "</body>\n</html>\n",
            ]
        )

    def to_text(self, show_image: bool = False) -> str:
        """
        Represent the slide page in text.

        Args:
            show_image (bool): Whether to show image captions.

        Returns:
            str: The text representation of the slide page.

        Raises:
            ValueError: If an image caption is not found.
        """
        text_content = "\n".join(
            [
                shape.text_frame.text.strip()
                for shape in self.shapes
                if shape.text_frame.is_textframe
            ]
        )
        if show_image:
            for image in self.shape_filter(Picture):
                if image.caption is None:
                    raise ValueError(
                        f"Caption not found for picture {image.shape_idx} of slide {image.slide_idx}"
                    )
                text_content += "\n" + "Image: " + image.caption
        return text_content

    @property
    def text_length(self) -> int:
        """
        Get the length of the text in the slide page.

        Returns:
            int: The length of the text.
        """
        return sum([len(shape.text_frame) for shape in self.shapes])

    def __iter__(self):
        """
        Iterate over all shapes in the slide page.

        Yields:
            ShapeElement: Each shape in the slide page.
        """
        for shape in self.shapes:
            if isinstance(shape, GroupShape):
                yield from shape
            else:
                yield shape

    def __len__(self) -> int:
        """
        Get the number of shapes in the slide page.

        Returns:
            int: The number of shapes.
        """
        return len(self.shapes)


class Presentation:
    """
    PPTAgent's representation of a presentation.
    Aiming at a more readable and editable interface.
    """

    def __init__(
        self,
        slides: List[SlidePage],
        error_history: List[Tuple[int, str]],
        slide_width: float,
        slide_height: float,
        file_path: str,
        num_pages: int,
    ) -> None:
        """
        Initialize the Presentation.

        Args:
            slides (List[SlidePage]): The slides in the presentation.
            error_history (List[Tuple[int, str]]): The error history.
            slide_width (float): The width of the slides.
            slide_height (float): The height of the slides.
            file_path (str): The path to the presentation file.
            num_pages (int): The number of pages in the presentation.
        """
        self.slides = slides
        self.error_history = error_history
        self.slide_width = slide_width
        self.slide_height = slide_height
        self.num_pages = num_pages
        self.source_file = file_path
        self.prs = PPTXPre(self.source_file)
        self.layout_mapping = {layout.name: layout for layout in self.prs.slide_layouts}
        self.prs.core_properties.last_modified_by = "PPTAgent"

    @classmethod
    def from_file(cls, file_path: str, config: Config) -> "Presentation":
        """
        Parse a Presentation from a file.

        Args:
            file_path (str): The path to the presentation file.
            config (Config): The configuration object.

        Returns:
            Presentation: The parsed Presentation.
        """
        prs = PPTXPre(file_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        slides = []
        error_history = []
        slide_idx = 0
        layouts = [layout.name for layout in prs.slide_layouts]
        num_pages = len(prs.slides)

        for slide in prs.slides:
            # Skip slides that won't be printed to PDF, as they are invisible
            if slide._element.get("show", 1) == "0":
                continue

            slide_idx += 1
            try:
                if slide.slide_layout.name not in layouts:
                    raise ValueError(
                        f"Slide layout {slide.slide_layout.name} not found"
                    )
                slides.append(
                    SlidePage.from_slide(
                        slide,
                        slide_idx - len(error_history),
                        slide_idx,
                        slide_width.pt,
                        slide_height.pt,
                        config,
                    )
                )
            except Exception as e:
                error_history.append((slide_idx, str(e)))
                if config.DEBUG:
                    logger.warning(
                        "Warning in slide %d of %s: %s",
                        slide_idx,
                        file_path,
                        traceback.format_exc(),
                    )

        return cls(
            slides, error_history, slide_width, slide_height, file_path, num_pages
        )

    def save(self, file_path: str, layout_only: bool = False) -> None:
        """
        Save the presentation to a file.

        Args:
            file_path (str): The path to save the presentation to.
            layout_only (bool): Whether to save only the layout.
        """
        self.clear_slides()
        for slide in self.slides:
            if layout_only:
                self.clear_images(slide.shapes)
            pptx_slide = self.build_slide(slide)
            if layout_only:
                self.clear_text(pptx_slide.shapes)
        self.prs.save(file_path)

    def build_slide(self, slide: SlidePage) -> PPTXSlide:
        """
        Build a slide in the presentation.
        """
        return slide.build(
            self.prs.slides.add_slide(self.layout_mapping[slide.slide_layout_name])
        )

    def clear_slides(self):
        """
        Delete all slides from the presentation.
        """
        while len(self.prs.slides) != 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]

    def clear_images(self, shapes: list[ShapeElement]):
        for shape in shapes:
            if isinstance(shape, GroupShape):
                self.clear_images(shape.data)
            elif isinstance(shape, Picture):
                shape.img_path = "resource/pic_placeholder.png"

    def clear_text(self, shapes: list[BaseShape]):
        for shape in shapes:
            if isinstance(shape, PPTXGroupShape):
                self.clear_text(shape.shapes)
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = "a" * len(run.text)

    def to_text(self, show_image: bool = False) -> str:
        """
        Represent the presentation in text.
        """
        return "\n----\n".join(
            [
                (
                    f"Slide {slide.slide_idx} of {len(self.prs.slides)}\n"
                    + (f"Title:{slide.slide_title}\n" if slide.slide_title else "")
                    + slide.to_text(show_image)
                )
                for slide in self.slides
            ]
        )

    def __len__(self) -> int:
        """
        Get the number of slides in the presentation.
        """
        return len(self.slides)
