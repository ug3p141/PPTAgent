# 1. 根据master 去完成生成内容到slide的映射关系
import pptx


class MasterGenerator:
    def __init__(self, prs: pptx.Presentation):
        self.prs = prs
        self.masters = self.prs.slide_layouts
        self.layout_slides_mapping = {
            # 母版的名字是不会重复的，所以可以用名字作为key
            i.name: [
                slide for slide in self.prs.slides if str(slide.slide_layout) == str(i)
            ]
            for i in self.masters
        }
        assert len(self.prs.slides) == sum(
            len(v) for v in self.layout_slides_mapping.values()
        )
        self.n_cluster = len(self.masters)

    def build_mapping(self):
        pass


if __name__ == "__main__":
    prs = pptx.Presentation("./resource/软件所党员组织关系专项整治工作部署.pptx")
    mg = MasterGenerator(prs)
    mg.build_mapping()
    for k, v in mg.layout_slides_mapping.items():
        print(k, "  :  ", len(v))
