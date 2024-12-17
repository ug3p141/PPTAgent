import inspect
import os
import re
import traceback
from copy import deepcopy
from dataclasses import dataclass
from enum import Enum
from functools import partial
from typing import Literal

import PIL
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.shapes.base import BaseShape
from pptx.util import Pt

from presentation import Closure, Picture, SlidePage
from utils import runs_merge


@dataclass
class HistoryMark:
    API_CALL_ERROR = "api_call_error"
    API_CALL_CORRECT = "api_call_correct"
    COMMENT_CORRECT = "comment_correct"
    COMMENT_ERROR = "comment_error"
    CODE_RUN_ERROR = "code_run_error"
    CODE_RUN_CORRECT = "code_run_correct"


class CodeExecutor:

    def __init__(self, retry_times: int):
        self.api_history = []
        self.command_history = []
        self.code_history = []
        self.retry_times = retry_times
        self.registered_functions = API_TYPES.all_funcs()
        self.function_regex = re.compile(r"^[a-z]+_[a-z_]+\(.+\)")

    def get_apis_docs(self, funcs: list[callable], show_example: bool = True):
        api_doc = []
        for func in funcs:
            sig = inspect.signature(func)
            params = []
            for name, param in sig.parameters.items():
                if name == "slide":
                    continue
                param_str = name
                if param.annotation != inspect.Parameter.empty:
                    param_str += f": {param.annotation.__name__}"
                if param.default != inspect.Parameter.empty:
                    param_str += f" = {repr(param.default)}"
                params.append(param_str)
            signature = f"def {func.__name__}({', '.join(params)})"
            if not show_example:
                api_doc.append(signature)
                continue
            doc = inspect.getdoc(func)
            if doc is not None:
                signature += f"\n\t{doc}"
            api_doc.append(signature)
        return "\n\n".join(api_doc)

    def execute_actions(
        self, actions: str, edit_slide: SlidePage, found_code: bool = False
    ):
        api_calls = actions.strip().split("\n")
        self.api_history.append(
            [HistoryMark.API_CALL_ERROR, edit_slide.slide_idx, actions]
        )
        for line_idx, line in enumerate(api_calls):
            try:
                if line_idx == len(api_calls) - 1 and not found_code:
                    raise ValueError(
                        "No code block found in the output, please output the api calls without any prefix."
                    )
                if line.startswith("def"):
                    raise PermissionError("The function definition were not allowed.")
                if line.startswith("#"):
                    if len(self.command_history) != 0:
                        self.command_history[-1][0] = HistoryMark.COMMENT_CORRECT
                    self.command_history.append([HistoryMark.COMMENT_ERROR, line, None])
                    continue
                if not self.function_regex.match(line):
                    continue
                found_code = True
                func = line.split("(")[0]
                if func not in self.registered_functions:
                    raise NameError(f"The function {func} is not defined.")
                # only one of clone and del can be used in a row
                if func.startswith("clone") or func.startswith("del"):
                    tag = func.split("_")[0]
                    if (
                        self.command_history[-1][-1] == None
                        or self.command_history[-1][-1] == tag
                    ):
                        self.command_history[-1][-1] = tag
                    else:
                        raise ValueError(
                            "Invalid command: Both 'clone_paragraph' and 'del_span'/'del_image' are used within a single command. "
                            "Each command must only perform one of these operations based on the quantity_change."
                        )
                self.code_history.append([HistoryMark.CODE_RUN_ERROR, line, None])
                partial_func = partial(self.registered_functions[func], edit_slide)
                eval(line, {}, {func: partial_func})
                self.code_history[-1][0] = HistoryMark.CODE_RUN_CORRECT
            except:
                trace_msg = traceback.format_exc()
                if len(self.code_history) != 0:
                    self.code_history[-1][-1] = trace_msg
                api_lines = (
                    "\n".join(api_calls[: line_idx - 1])
                    + f"\n--> Error Line: {line}\n"
                    + "\n".join(api_calls[line_idx + 1 :])
                )
                return api_lines, trace_msg
        if len(self.command_history) != 0:
            self.command_history[-1][0] = HistoryMark.COMMENT_CORRECT
        self.api_history[-1][0] = HistoryMark.API_CALL_CORRECT


# supporting functions
def element_index(slide: SlidePage, element_id: int):
    for shape in slide:
        if shape.shape_idx == element_id:
            return shape
    raise IndexError(f"Cannot find element {element_id}, is it deleted or not exist?")


def replace_run(paragraph_id: int, run_id: int, new_text: str, shape: BaseShape):
    para = shape.text_frame.paragraphs[paragraph_id]
    run = runs_merge(para)[run_id]
    run.text = new_text


def clone_para(paragraph_id: int, shape: BaseShape):
    para = shape.text_frame.paragraphs[paragraph_id]
    shape.text_frame.paragraphs[-1]._element.addnext(parse_xml(para._element.xml))


def del_run(paragraph_id: int, run_id: int, shape: BaseShape):
    para = shape.text_frame.paragraphs[paragraph_id]
    run = runs_merge(para)[run_id]
    run._r.getparent().remove(run._r)
    if len(para.runs) == 0:
        para._element.getparent().remove(para._element)


HORIZONTAL_ALIGN = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}
VERTICAL_ALIGN = {
    "top": MSO_ANCHOR.TOP,
    "center": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}
DEFAULT_FONT_SIZE = Pt(12)


def set_font(
    bold: bool,
    font_size_delta: int | None,
    horizontal_align: Literal["left", "center", "right", None],
    vertical_align: Literal["top", "center", "bottom", None],
    text_shape: BaseShape | None,
):
    horizontal_align = HORIZONTAL_ALIGN.get(horizontal_align, None)
    vertical_align = VERTICAL_ALIGN.get(vertical_align, None)
    if bold is not None:
        text_shape.text_frame.font.bold = bold

    textframe_font_size = text_shape.text_frame.font.size or Pt(12)

    for paragraph in text_shape.text_frame.paragraphs:
        if horizontal_align is not None:
            paragraph.alignment = horizontal_align
        if vertical_align is not None:
            paragraph.vertical_anchor = vertical_align
        if font_size_delta is None:
            continue
        para_font_size = (paragraph.font.size or textframe_font_size) + font_size_delta
        for run in paragraph.runs:
            if run.font.size is not None:
                run.font.size = run.font.size + font_size_delta
        paragraph.font.size = para_font_size


def set_size(
    left_delta: int,
    right_delta: int,
    top_delta: int,
    bottom_delta: int,
    shape: BaseShape,
):
    if left_delta is not None:
        shape.left = Pt(shape.left.pt + left_delta)
    if right_delta is not None:
        shape.width = Pt(shape.width.pt + right_delta)
    if bottom_delta is not None:
        shape.height = Pt(shape.height.pt + bottom_delta)
    if top_delta is not None:
        shape.top = Pt(shape.top.pt + top_delta)


# api functions
def del_span(slide: SlidePage, div_id: int, paragraph_id: int, span_id: int):
    shape = element_index(slide, div_id)
    assert (
        shape.text_frame.is_textframe
    ), "The element does not have a text frame, please check the element id and type of element."
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.idx == paragraph_id:
            para = paragraph
            break
    else:
        raise IndexError(
            f"Cannot find the paragraph {paragraph_id} of the element {div_id},"
            "may refer to a non-existed paragraph or you haven't cloned enough paragraphs beforehand."
        )
    for run in para.runs:
        if run.idx == span_id:
            para.runs.remove(run)
            shape._closures["delete"].append(
                Closure(
                    partial(del_run, para.real_idx, span_id), para.real_idx, span_id
                )
            )
            return
    raise IndexError(
        f"Cannot find the span {span_id} in the paragraph {paragraph_id} of the element {div_id}, may refer to a non-existent span."
    )


def del_image(slide: SlidePage, figure_id: int):
    shape = element_index(slide, figure_id)
    assert isinstance(shape, Picture), "The element is not a Picture."
    slide.shapes.remove(shape)


def replace_span(
    slide: SlidePage, div_id: int, paragraph_id: int, span_id: int, text: str
):
    shape = element_index(slide, div_id)
    assert (
        shape.text_frame.is_textframe
    ), "The element does not have a text frame, please check the element id and type of element."
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.idx == paragraph_id:
            para = paragraph
            break
    else:
        raise IndexError(
            f"Cannot find the paragraph {paragraph_id} of the element {div_id},"
            "may refer to a non-existed paragraph or you haven't cloned enough paragraphs beforehand."
        )
    for run in para.runs:
        if run.idx == span_id:
            run.text = text
            shape._closures["replace"].append(
                Closure(
                    partial(replace_run, para.real_idx, span_id, text),
                    para.real_idx,
                    span_id,
                )
            )
            return
    raise IndexError(
        f"Cannot find the span {span_id} in the paragraph {paragraph_id} of the element {div_id},"
        "Please: "
        "1. check if you refer to a non-existed span."
        "2. check if you already deleted it, ensure to remove span elements from the end of the paragraph first"
        "3. consider merging adjacent replace_span operations."
    )


def replace_image(slide: SlidePage, img_id: int, image_path: str):
    if not os.path.exists(image_path):
        raise ValueError(
            f"The image {image_path} does not exist, consider use del_image if image_path in the given command is faked"
        )
    shape = element_index(slide, img_id)
    assert isinstance(shape, Picture), "The element is not a Picture."
    img_size = PIL.Image.open(image_path).size
    r = min(shape.width / img_size[0], shape.height / img_size[1])
    new_width = int(img_size[0] * r)
    new_height = int(img_size[1] * r)
    shape.width = Pt(new_width)
    shape.height = Pt(new_height)
    shape.img_path = image_path


def clone_paragraph(slide: SlidePage, div_id: int, paragraph_id: int):
    # The cloned paragraph will have a paragraph_id one greater than the current maximum in the parent element.
    shape = element_index(slide, div_id)
    assert (
        shape.text_frame.is_textframe
    ), "The element does not have a text frame, please check the element id and type of element."
    max_idx = max([para.idx for para in shape.text_frame.paragraphs])
    for para in shape.text_frame.paragraphs:
        if para.idx != paragraph_id:
            continue
        shape.text_frame.paragraphs.append(deepcopy(para))
        shape.text_frame.paragraphs[-1].idx = max_idx + 1
        shape.text_frame.paragraphs[-1].real_idx = len(shape.text_frame.paragraphs) - 1
        shape._closures["clone"].append(
            Closure(
                partial(clone_para, para.real_idx),
                para.real_idx,
            )
        )
        return
    raise IndexError(
        f"Cannot find the paragraph {paragraph_id} of the element {div_id}, may refer to a non-existed paragraph."
    )


class API_TYPES(Enum):
    Agent = [
        del_span,
        del_image,
        clone_paragraph,
        replace_span,
        replace_image,
    ]

    @classmethod
    def all_funcs(cls) -> dict[str, callable]:
        funcs = {}
        for attr in dir(cls):
            if attr.startswith("__"):
                continue
            funcs |= {func.__name__: func for func in getattr(cls, attr).value}
        return funcs


if __name__ == "__main__":
    print(CodeExecutor(0).get_apis_docs(API_TYPES.Agent.value))
