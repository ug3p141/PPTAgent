import json
import os
from collections import defaultdict
from copy import deepcopy
from datetime import datetime

import json_repair
import torch
from jinja2 import Template
from pptx import Presentation as PPTXPre
from torch import cosine_similarity

import apis
from apis import API_TYPES, code_executor
from llms import agent_model
from model_utils import get_text_embedding
from presentation import Presentation
from utils import app_config, clear_slides, pexists, pjoin, print


class PPTAgent:
    def __init__(
        self,
        presentation: Presentation,
        template: dict,
        images: dict[str, str],
        num_slides: int,
        doc_json: dict[str, str],
        functional_keys: list[str],
    ):
        self.presentation = presentation
        self.slide_templates = template
        self.doc_json = doc_json
        self.num_slides = num_slides
        self.image_information = "\n".join(
            [
                f"Image path: {k}, size: {v[1][0]}*{v[1][1]} px\n caption: {v[0]}"
                for k, v in images.items()
            ]
        )
        self.functional_keys = functional_keys

        apis.image_stats = images
        apis.image_usage = defaultdict(int)

        self.gen_prs = deepcopy(presentation)
        self.gen_prs.slides = []

        output_dir = pjoin(app_config.RUN_DIR, "agent")
        os.makedirs(output_dir, exist_ok=True)
        self.outline_file = pjoin(output_dir, "slide_outline.json")
        self.cache_file = pjoin(output_dir, "generating_steps.json")

    # TODO vision edit的结果单独保存
    def work(self):
        if pexists(self.outline_file):
            self.outline = json.load(open(self.outline_file, "r"))
        else:
            self.outline = json_repair.loads(self.generate_outline())
            json.dump(
                self.outline, open(self.outline_file, "w"), ensure_ascii=False, indent=4
            )
        meta_data = "\n".join(
            [f"{k}: {v}" for k, v in self.doc_json["metadata"].items()]
        )
        self.metadata = f"\nMetadatao of Presentation: \n{meta_data}\nCurrent Time: {datetime.now().strftime('%Y-%m-%d')}\n"
        self.simple_outline = "Outline of Presentation: \n" + "\n".join(
            [
                f"Slide {slide_idx+1}: {slide_title}"
                for slide_idx, slide_title in enumerate(self.outline)
            ]
        )
        self.generate_slides()

    def generate_slides(self, use_cache=True):
        edit_template = Template(open("prompts/agent/edit.txt").read())

        self.gen_prs.slides = []
        steps = []
        layout_names = list(self.slide_templates.keys())
        layout_embeddings = torch.stack(get_text_embedding(layout_names, 32))

        if use_cache and pexists(self.cache_file):
            steps, code_executor.api_history, code_executor.code_history = json.load(
                open(self.cache_file, "r")
            )
            for step in steps:
                apis.slide = deepcopy(self.presentation.slides[step[0]])
                code_executor.execute_apis(step[2], step[3])
                self.gen_prs.slides.append(apis.slide)
        for slide_idx, (slide_title, slide) in enumerate(self.outline.items()):
            if use_cache and len(steps) != slide_idx:
                continue
            images = "No Images"
            if any(
                [
                    i in slide["layout"]
                    for i in ["picture", "chart", "table", "diagram", "freeform"]
                ]
            ):
                images = self.image_information
            slide_content = self.get_slide_content(slide_idx, slide_title, slide)
            if slide["layout"] not in self.slide_templates:
                layout_sim = torch.cosine_similarity(
                    get_text_embedding(slide["layout"], 1)[0], layout_embeddings
                )
                slide["layout"] = layout_names[layout_sim.argmax().item()]
            template_id = (
                max(
                    self.slide_templates[slide["layout"]],
                    key=lambda x: len(self.presentation.slides[x - 1].shapes),
                )
                - 1
            )
            apis.slide = deepcopy(self.presentation.slides[template_id])
            edit_prompt = edit_template.render(
                api_documentation=code_executor.get_apis_docs(
                    [API_TYPES.TEXT_EDITING, API_TYPES.IMAGE_EDITING]
                ),
                edit_target=apis.slide.to_html(),
                content=self.simple_outline + self.metadata + slide_content,
                images=images,
            )
            code_executor.execute_apis(
                edit_prompt,
                apis=agent_model(
                    edit_prompt,
                ),
                # 是否在这步提供图片
            )
            steps.append((template_id, *code_executor.api_history[-1]))
            self.gen_prs.slides.append(apis.slide)
            if use_cache:
                json.dump(
                    (steps, code_executor.api_history, code_executor.code_history),
                    open(self.cache_file, "w"),
                )
        self.gen_prs.save(pjoin(app_config.RUN_DIR, "agent", "final.pptx"))

    def generate_outline(self):
        template = Template(open("prompts/agent/outline.txt").read())
        prompt = template.render(
            num_slides=self.num_slides,
            layouts="\n".join(
                set(self.slide_templates.keys()).difference(self.functional_keys)
            ),
            functional_keys="\n".join(self.functional_keys),
            json_content=self.doc_json,
            images=self.image_information,
        )
        return agent_model(prompt)

    def get_slide_content(self, slide_idx: int, slide_title: str, slide: dict):
        slide_content = f"Slide-{slide_idx+1} Title: {slide_title}\nSlide Description: {slide['description']}\n"
        if len(slide["subsection_keys"]) != 0:
            slide_content += "Slide Reference Text: "
            for key in slide["subsection_keys"]:
                for section in self.doc_json["sections"]:
                    for subsection in section["subsections"]:
                        if key in subsection:
                            slide_content += f"SubSection {key}: {subsection[key]}\n"
        return slide_content
