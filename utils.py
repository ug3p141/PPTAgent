import json
import os
import shutil
import subprocess
import tempfile
import traceback
import xml.etree.ElementTree as ET
from types import SimpleNamespace

import json_repair
import requests
from lxml import etree
from pdf2image import convert_from_path
from pptx.dml.fill import _NoFill, _NoneFill
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.util import Length
from rich import print
from tenacity import RetryCallState, retry, stop_after_attempt, wait_fixed

IMAGE_EXTENSIONS = {"bmp", "jpg", "jpeg", "pgm", "png", "ppm", "tif", "tiff", "webp"}


def get_slide_content(doc_json: dict, slide_title: str, slide: dict):
    slide_desc = slide.get("description", "")
    slide_content = f"Title: {slide_title}\nSlide Description: {slide_desc}\n"
    if len(slide.get("subsection_keys", [])) != 0:
        slide_content += "Slide Reference Text: "
        for key in slide["subsection_keys"]:
            for section in doc_json["sections"]:
                for subsection in section.get("subsections", []):
                    if key in subsection:
                        slide_content += f"SubSection {key}: {subsection[key]}\n"
    return slide_content


def tenacity_log(retry_state: RetryCallState):
    print(f"Retry attempt {retry_state.attempt_number}")
    traceback.print_tb(retry_state.outcome.exception().__traceback__)


def get_json_from_response(response: str):
    l, r = response.rfind("```json"), response.rfind("```")
    try:
        if l == -1 or r == -1:
            return json_repair.loads(response)
        return json_repair.loads(response[l + 7 : r].strip())
    except:
        raise RuntimeError("Failed to parse JSON from response")

tenacity = retry(
    wait=wait_fixed(3), stop=stop_after_attempt(3), after=tenacity_log, reraise=True
)

def parse_pdf(file: str, output_dir: str, api: str):
    os.makedirs(output_dir, exist_ok=True)
    with tempfile.NamedTemporaryFile(suffix=".zip") as temp_zip:
        with open(file, "rb") as f:
            response = requests.post(api, files={"pdf": f})
            response.raise_for_status()
            temp_zip.write(response.content)
            temp_zip.flush()

        shutil.unpack_archive(temp_zip.name, output_dir)

def ppt_to_images(file: str, output_dir: str):
    os.makedirs(output_dir, exist_ok=True)
    with tempfile.TemporaryDirectory() as temp_dir:
        command_list = [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            file,
            "--outdir",
            temp_dir,
        ]
        subprocess.run(command_list, check=True, stdout=subprocess.DEVNULL)

        for f in os.listdir(temp_dir):
            if not f.endswith(".pdf"):
                continue
            temp_pdf = pjoin(temp_dir, f)
            images = convert_from_path(temp_pdf, dpi=72)
            for i, img in enumerate(images):
                img.save(pjoin(output_dir, f"slide_{i+1:04d}.jpg"))
            return

        raise RuntimeError("No PDF file was created in the temporary directory")

def filename_normalize(filename: str):
    return filename.replace("/", "_").replace(" ", "_").replace("\\", "_")

def set_proxy(proxy_url: str):
    os.environ["http_proxy"] = proxy_url
    os.environ["https_proxy"] = proxy_url

def get_text_inlinestyle(para: dict, stylish: bool):
    if not stylish:
        return ""
    font = SimpleNamespace(**para["font"])
    font_size = f"font-size: {Length(font.size).pt}pt;" if font.size else ""
    # font_family = f"font-family: {font.name};" if font.name else ""
    font_color = f"color='{font.color}';" if font.color else ""
    font_bold = "font-weight: bold;" if font.bold else ""
    return 'style="{}"'.format("".join([font_size, font_color, font_bold]))


def extract_fill(shape: BaseShape):
    if "fill" not in dir(shape):
        return None
    fill_dict = {
        "fill_xml": shape.fill._xPr.xml,
    } | {k: v for k, v in object_to_dict(shape.fill).items() if "color" in k}
    if not isinstance(shape.fill._fill, (_NoneFill, _NoFill)):
        fill_dict["type"] = shape.fill.type.name.lower()
    return fill_dict


def apply_fill(shape: BaseShape, fill: dict):
    if fill is None:
        return
    new_element = etree.fromstring(fill["fill_xml"])
    shape.fill._xPr.getparent().replace(shape.fill._xPr, new_element)


def parse_groupshape(groupshape: GroupShape):
    assert isinstance(groupshape, GroupShape)
    group_top_left_x = groupshape.left
    group_top_left_y = groupshape.top
    group_width = groupshape.width
    group_height = groupshape.height
    shape_top_left_x = min([sp.left for sp in groupshape.shapes])
    shape_top_left_y = min([sp.top for sp in groupshape.shapes])
    shape_width = (
        max([sp.left + sp.width for sp in groupshape.shapes]) - shape_top_left_x
    )
    shape_height = (
        max([sp.top + sp.height for sp in groupshape.shapes]) - shape_top_left_y
    )
    group_shape_xy = []
    for sp in groupshape.shapes:
        group_shape_left = (
            sp.left - shape_top_left_x
        ) * group_width / shape_width + group_top_left_x
        group_shape_top = (
            sp.top - shape_top_left_y
        ) * group_height / shape_height + group_top_left_y
        group_shape_width = sp.width * group_width / shape_width
        group_shape_height = sp.height * group_height / shape_height
        group_shape_xy.append(
            {
                "left": Length(group_shape_left),
                "top": Length(group_shape_top),
                "width": Length(group_shape_width),
                "height": Length(group_shape_height),
            }
        )
    return group_shape_xy


def is_primitive(obj):
    """
    判断对象或该集合包含的所有对象是否是基本类型。

    参数:
    obj: 要判断的对象

    返回:
    如果对象是基本类型，返回True，否则返回False
    """
    if isinstance(obj, (list, tuple, set, frozenset)):
        return all(is_primitive(item) for item in obj)
    return isinstance(
        obj, (int, float, complex, bool, str, bytes, bytearray, type(None))
    )


DEFAULT_EXCLUDE = set(["element", "language_id", "ln", "placeholder_format"])


def object_to_dict(obj, result=None, exclude=None):
    """
    将对象的非隐藏属性拷贝到一个字典中。

    参数:
    obj: 要拷贝属性的对象

    返回:
    包含对象非隐藏属性的字典
    """
    if result is None:
        result = {}
    exclude = DEFAULT_EXCLUDE.union(exclude or set())
    for attr in dir(obj):
        if attr in exclude:
            continue
        try:
            if not attr.startswith("_") and not callable(getattr(obj, attr)):
                attr_value = getattr(obj, attr)
                if "real" in dir(attr_value):
                    attr_value = attr_value.real
                if is_primitive(attr_value):
                    result[attr] = attr_value
        except:
            pass
    return result


def dict_to_object(dict: dict, obj: object, exclude=None):
    """
    从字典中恢复对象的属性。

    参数:
    d: 包含对象属性的字典
    obj: 要恢复属性的对象

    返回:
    恢复属性后的对象
    """
    if exclude is None:
        exclude = set()
    for key, value in dict.items():
        if key not in exclude:
            setattr(obj, key, value)


class Config:
    def __init__(self, session_id=None, rundir=None, debug=True):
        self.DEBUG = debug
        if session_id is not None:
            self.set_session(session_id)
        if rundir is not None:
            self.set_rundir(rundir)

    def set_session(self, session_id):
        self.session_id = session_id
        self.set_rundir(f"./runs/{session_id}")

    def set_rundir(self, rundir: str):
        self.RUN_DIR = rundir
        self.IMAGE_DIR = pjoin(self.RUN_DIR, "images")
        for the_dir in [self.RUN_DIR, self.IMAGE_DIR]:
            os.makedirs(the_dir, exist_ok=True)

    def set_debug(self, debug: bool):
        self.DEBUG = debug

    def remove_rundir(self):
        if pexists(self.RUN_DIR):
            shutil.rmtree(self.RUN_DIR)
        if pexists(self.IMAGE_DIR):
            shutil.rmtree(self.IMAGE_DIR)


pjoin = os.path.join
pexists = os.path.exists
pbasename = os.path.basename

if __name__ == "__main__":
    config = Config()
    print(config)
