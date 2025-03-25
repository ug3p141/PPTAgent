from pptagent.apis import CodeExecutor, API_TYPES, replace_para
from pptagent.utils import package_join
from pptx import Presentation
from test.conftest import test_config


def test_api_docs():
    executor = CodeExecutor(3)
    docs = executor.get_apis_docs(API_TYPES.Agent.value)
    assert len(docs) > 0


def test_replace_para():
    text = "这是一个**加粗和*斜体*文本**，还有*斜体和`Code def a+b`*，~~删除~~，[链接](http://example.com)"
    prs = Presentation(package_join(test_config.ppt))
    slide = prs.slides[0]
    replace_para(0, text, slide.shapes[0])
    runs = slide.shapes[0].text_frame.paragraphs[0].runs
    assert runs[1].font.bold
    assert runs[2].font.bold and runs[2].font.italic
    assert runs[6].font.name == "Consolas"
    assert runs[8].font.strikethrough
    assert runs[10].hyperlink.address == "http://example.com"
