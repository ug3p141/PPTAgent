import re
from dataclasses import dataclass
from typing import Callable, Dict, List, Optional, Type, TypeVar, Union

from pptx.chart.chart import Chart as PPTXChart
from pptx.slide import Slide as PPTXSlide
from pptx.shapes.connector import Connector as PPTXConnector
from pptx.oxml import parse_xml
from pptx.shapes.autoshape import Shape as PPTXAutoShape
from pptx.shapes.graphfrm import GraphicFrame as PPTXGraphicalFrame
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape as PPTXGroupShape
from pptx.shapes.picture import Picture as PPTXPicture
from pptx.shapes.placeholder import PlaceholderPicture, SlidePlaceholder
from pptx.text.text import _Paragraph
from pptx.table import Table as PPTXTable
from pptx.enum.shapes import MSO_SHAPE_TYPE


from utils import (
    IMAGE_EXTENSIONS,
    Config,
    apply_fill,
    dict_to_object,
    extract_fill,
    get_font_style,
    merge_dict,
    object_to_dict,
    parse_groupshape,
    pexists,
    pjoin,
    runs_merge,
    wmf_to_images,
)

INDENT = "\t"
T = TypeVar("T", bound="ShapeElement")


@dataclass
class StyleArg:
    """
    A class to represent style arguments for HTML conversion.
    """

    paragraph_id: bool = True
    element_id: bool = True
    font_style: bool = True
    # todo 这里还没实现
    fill_style: bool = True
    area: bool = False
    size: bool = False
    geometry: bool = False
    show_name: bool = False
    show_image: bool = True
    show_content: bool = True
    show_semantic_name: bool = False

    @classmethod
    def all_true(cls) -> "StyleArg":
        """
        Create a StyleArg instance with all options enabled.

        Returns:
            StyleArg: A StyleArg instance with all options enabled.
        """
        return cls(
            area=True,
            size=True,
            geometry=True,
            show_semantic_name=True,
        )


@dataclass
class Closure:
    """
    A class to represent a closure that can be applied to a shape.
    """

    closure: Callable
    paragraph_id: int = -1

    def apply(self, shape: BaseShape) -> None:
        """
        Apply the closure to a shape.

        Args:
            shape (BaseShape): The shape to apply the closure to.
        """
        self.closure(shape)

    def __gt__(self, other: "Closure") -> bool:
        """
        Compare closures based on paragraph_id.

        Args:
            other (Closure): Another closure to compare with.

        Returns:
            bool: True if this closure's paragraph_id is greater than the other's.
        """
        if self.paragraph_id != other.paragraph_id:
            return self.paragraph_id > other.paragraph_id


class Paragraph:
    """
    A class to represent a paragraph in a text frame.
    """

    def __init__(self, paragraph: _Paragraph, idx: int):
        """
        Initialize a Paragraph.

        Args:
            paragraph (_Paragraph): The paragraph object.
            idx (int): The index of the paragraph.
        """
        run = runs_merge(paragraph)
        self.idx = idx
        self.real_idx = idx
        self.bullet = paragraph.bullet
        if run is None:
            self.idx = -1
            return
        self.font = merge_dict(
            object_to_dict(paragraph.font), [object_to_dict(run.font)]
        )
        self.text = re.sub(r"(_x000B_|\\x0b)", " ", paragraph.text)

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the paragraph to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the paragraph.

        Raises:
            ValueError: If the paragraph is not valid.
        """
        if self.idx == -1:
            raise ValueError(f"paragraph {self.idx} is not valid")
        tag = "li" if self.bullet else "p"
        id_str = f" id='{self.idx}'" if style_args.paragraph_id else ""
        font_style = get_font_style(self.font)
        style_str = (
            f" style='{font_style}'" if style_args.font_style and font_style else ""
        )
        if self.bullet:
            style_str += f" bullet-type='{self.bullet}'"
        return f"<{tag}{id_str}{style_str}>{self.text}</{tag}>"

    def __repr__(self) -> str:
        """
        Get a string representation of the paragraph.

        Returns:
            str: A string representation of the paragraph.
        """
        return f"Paragraph-{self.idx}: {self.text}"


class TextFrame:
    """
    A class to represent a text frame in a shape.
    """

    def __init__(self, shape: BaseShape, level: int):
        """
        Initialize a TextFrame.

        Args:
            shape (BaseShape): The shape containing the text frame.
            level (int): The indentation level.
        """
        if not shape.has_text_frame:
            self.is_textframe = False
            return
        self.paragraphs = [
            Paragraph(paragraph, idx)
            for idx, paragraph in enumerate(shape.text_frame.paragraphs)
        ]
        para_offset = 0
        for para in self.paragraphs:
            if para.idx == -1:
                para_offset += 1
            else:
                para.idx = para.idx - para_offset
        if len(self.paragraphs) == 0:
            self.is_textframe = False
            return
        self.level = level
        self.text = shape.text
        self.is_textframe = True
        self.font = merge_dict(
            object_to_dict(shape.text_frame.font),
            [para.font for para in self.paragraphs if para.idx != -1],
        )

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the text frame to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the text frame.
        """
        if not self.is_textframe:
            return ""
        repr_list = [
            para.to_html(style_args) for para in self.paragraphs if para.idx != -1
        ]
        return "\n".join([INDENT * self.level + repr for repr in repr_list])

    def __repr__(self) -> str:
        """
        Get a string representation of the text frame.

        Returns:
            str: A string representation of the text frame.
        """
        if not self.is_textframe:
            return "TextFrame: null"
        return f"TextFrame: {self.paragraphs}"

    def __len__(self) -> int:
        """
        Get the length of the text in the text frame.

        Returns:
            int: The length of the text.
        """
        if not self.is_textframe:
            return 0
        return len(self.text)


class Background:
    """
    A class to represent a slide background.
    """

    def __init__(self, slide: PPTXSlide):
        """
        Initialize a Background.

        Args:
            slide (PPTXSlide): The slide containing the background.
        """
        background = slide.background
        self.xml = background._element.xml

    def build(self, slide: PPTXSlide) -> None:
        """
        Build the background in a slide.

        Args:
            slide (PPTXSlide): The slide to build the background in.
        """
        pass

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the background to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the background.
        """
        pass


class ShapeElement:
    """
    Base class for shape elements in a presentation.
    """

    def __init__(
        self,
        slide_idx: int,
        shape_idx: int,
        style: Dict,
        data: List,
        text_frame: TextFrame,
        slide_area: float,
        level: int,
    ):
        """
        Initialize a ShapeElement.

        Args:
            slide_idx (int): The index of the slide.
            shape_idx (int): The index of the shape.
            style (Dict): The style of the shape.
            data (List): The data of the shape.
            text_frame (TextFrame): The text frame of the shape.
            slide_area (float): The area of the slide.
            level (int): The indentation level.
        """
        self.slide_idx = slide_idx
        self.shape_idx = shape_idx
        self.style = style
        self.data = data
        self.text_frame = text_frame
        self._closure_keys = ["clone", "replace", "delete", "style"]
        self._closures: Dict[str, List[Closure]] = {
            key: [] for key in self._closure_keys
        }
        self.slide_area = slide_area
        self.level = level
        self.xml = None  # Will be set in from_shape

    @classmethod
    def from_shape(
        cls: Type[T],
        slide_idx: int,
        shape_idx: int,
        shape: BaseShape,
        config: Config,
        slide_area: float,
        level: int = 0,
    ) -> T:
        """
        Create a ShapeElement from a BaseShape.

        Args:
            slide_idx (int): The index of the slide.
            shape_idx (int): The index of the shape.
            shape (BaseShape): The shape object.
            config (Config): The configuration object.
            slide_area (float): The area of the slide.
            level (int): The indentation level.

        Returns:
            T: The created ShapeElement.

        Raises:
            ValueError: If nested group shapes are not allowed.
        """
        if shape_idx > 100 and isinstance(shape, PPTXGroupShape):
            raise ValueError("Nested group shapes are not allowed")

        # Extract line properties if available
        line = None
        if hasattr(shape, "line") and shape.line._ln is not None:
            line = {
                "fill_xml": extract_fill(shape.line),
                "width": shape.line.width,
                "dash_style": shape.line.dash_style,
            }

        # Extract fill properties
        fill = extract_fill(shape)

        # Create style dictionary
        style = {
            "shape_bounds": {
                "width": shape.width,
                "height": shape.height,
                "left": shape.left,
                "top": shape.top,
            },
            "shape_type": str(shape.shape_type).split("(")[0].lower(),
            "rotation": shape.rotation,
            "fill_xml": fill,
            "line": line,
            "name": shape.name,
        }

        # Determine semantic name
        try:
            # For auto shapes (rectangle, oval, triangle, star...)
            autoshape = shape.auto_shape_type
            assert autoshape is not None
            style["semantic_name"] = str(autoshape).split()[0].lower().strip()
        except:
            # For other shapes (freeform, connector, table, chart...)
            style["semantic_name"] = str(shape.shape_type).split("(")[0].lower().strip()

        # Create text frame
        text_frame = TextFrame(shape, level + 1)

        # Create appropriate shape element based on shape type
        shape_class = SHAPECAST.get(shape.shape_type, UnsupportedShape)
        obj = shape_class.from_shape(
            slide_idx,
            shape_idx,
            shape,
            style,
            text_frame,
            config,
            slide_area,
            level,
        )

        # Store XML for later use
        obj.xml = shape._element.xml
        # ? This is for debug use, mask to enable pickling
        # obj.shape = shape
        return obj

    def build(self, slide: PPTXSlide) -> BaseShape:
        """
        Build the shape element in a slide.

        Args:
            slide (PPTXSlide): The slide to build the shape in.

        Returns:
            BaseShape: The built shape.
        """
        return slide.shapes._shape_factory(
            slide.shapes._spTree.insert_element_before(parse_xml(self.xml), "p:extLst")
        )

    def __repr__(self) -> str:
        """
        Get a string representation of the shape element.

        Returns:
            str: A string representation of the shape element.
        """
        return f"{self.__class__.__name__}: shape {self.shape_idx} of slide {self.slide_idx}"

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the shape element to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the shape element.

        Raises:
            NotImplementedError: If not implemented in a subclass.
        """
        raise NotImplementedError(
            f"to_html not implemented for {self.__class__.__name__}"
        )

    @property
    def closures(self) -> List[Closure]:
        """
        Get the closures associated with the shape element.

        Returns:
            List[Closure]: A list of closures.
        """
        closures = []
        closures.extend(sorted(self._closures["clone"]))
        closures.extend(self._closures["replace"] + self._closures["style"])
        closures.extend(sorted(self._closures["delete"], reverse=True))
        return closures

    @property
    def indent(self) -> str:
        """
        Get the indentation string for the shape element.

        Returns:
            str: The indentation string.
        """
        return "\t" * self.level

    @property
    def left(self) -> float:
        """
        Get the left position of the shape element.

        Returns:
            float: The left position in points.
        """
        return self.style["shape_bounds"]["left"].pt

    @left.setter
    def left(self, value: float) -> None:
        """
        Set the left position of the shape element.

        Args:
            value (float): The left position in points.
        """
        self.style["shape_bounds"]["left"] = value

    @property
    def top(self) -> float:
        """
        Get the top position of the shape element.

        Returns:
            float: The top position in points.
        """
        return self.style["shape_bounds"]["top"].pt

    @top.setter
    def top(self, value: float) -> None:
        """
        Set the top position of the shape element.

        Args:
            value (float): The top position in points.
        """
        self.style["shape_bounds"]["top"] = value

    @property
    def width(self) -> float:
        """
        Get the width of the shape element.

        Returns:
            float: The width in points.
        """
        return self.style["shape_bounds"]["width"].pt

    @width.setter
    def width(self, value: float) -> None:
        """
        Set the width of the shape element.

        Args:
            value (float): The width in points.
        """
        self.style["shape_bounds"]["width"] = value

    @property
    def height(self) -> float:
        """
        Get the height of the shape element.

        Returns:
            float: The height in points.
        """
        return self.style["shape_bounds"]["height"].pt

    @height.setter
    def height(self, value: float) -> None:
        """
        Set the height of the shape element.

        Args:
            value (float): The height in points.
        """
        self.style["shape_bounds"]["height"] = value

    @property
    def area(self) -> float:
        """
        Get the area of the shape element.

        Returns:
            float: The area in square points.
        """
        return self.width * self.height

    @property
    def semantic_name(self) -> Optional[str]:
        """
        Get the semantic name of the shape element.

        Returns:
            Optional[str]: The semantic name, or None if not set.
        """
        return self.style.get("semantic_name", None)

    @semantic_name.setter
    def semantic_name(self, value: str) -> None:
        """
        Set the semantic name of the shape element.

        Args:
            value (str): The semantic name.
        """
        self.style["semantic_name"] = value

    def get_inline_style(self, style_args: StyleArg) -> str:
        """
        Get the inline style for the shape element.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The inline style string.
        """
        id_str = f" id='{self.shape_idx}'" if style_args.element_id else ""
        data_attrs = []
        styles = []

        # Add data attributes
        if style_args.area:
            data_attrs.append(
                f"data-relative-area={self.area*100/self.slide_area:.2f}%;"
            )
        if style_args.show_name:
            data_attrs.append(f"data-shapeName='{self.style['name']}'")
        if style_args.show_semantic_name and self.semantic_name is not None:
            data_attrs.append(f"data-semanticName='{self.semantic_name}'")

        # Add style attributes
        if style_args.size:
            styles.append(f"width: {self.width}pt; height: {self.height}pt;")
        if style_args.geometry:
            styles.append(f"left: {self.left}pt; top: {self.top}pt;")
        if style_args.font_style and self.text_frame.is_textframe:
            font_style = get_font_style(self.text_frame.font)
            if font_style:
                styles.append(font_style)

        # Combine attributes
        if len(styles) != 0:
            id_str += " style='" + " ".join(styles) + "'"
        if len(data_attrs) != 0:
            id_str += " " + " ".join(data_attrs)

        return id_str


class UnsupportedShape(ShapeElement):
    """
    A class to represent an unsupported shape.
    """

    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: BaseShape,
        *args,
        **kwargs,
    ) -> None:
        """
        Create an UnsupportedShape from a BaseShape.

        Raises:
            ValueError: Always, as the shape is unsupported.
        """
        raise ValueError(f"Unsupported shape {shape.shape_type}")


class TextBox(ShapeElement):
    """
    A class to represent a text box shape.
    """

    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: TextFrame,
        style: Dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ) -> "TextBox":
        """
        Create a TextBox from a TextFrame.

        Returns:
            TextBox: The created TextBox.
        """
        return cls(slide_idx, shape_idx, style, [], text_frame, slide_area, level)

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the text box to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the text box.
        """
        content = self.text_frame.to_html(style_args)
        if not style_args.show_content:
            content = ""
        return (
            f"{self.indent}<div{self.get_inline_style(style_args)}>\n"
            + content
            + f"\n{self.indent}</div>\n"
        )


class Picture(ShapeElement):
    """
    A class to represent a picture shape.
    """

    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: PPTXPicture,
        style: Dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ) -> "Picture":
        """
        Create a Picture from a PPTXPicture.

        Returns:
            Picture: The created Picture.

        Raises:
            ValueError: If the image type is unsupported.
        """
        img_path = pjoin(
            config.IMAGE_DIR,
            f"{shape.image.sha1}.{shape.image.ext}",
        )

        # Handle WMF images
        if shape.image.ext == "wmf":
            img_path = img_path.replace(".wmf", ".jpg")
            if not pexists(img_path):
                wmf_to_images(shape.image.blob, img_path)
        # Check for supported image types
        elif shape.image.ext not in IMAGE_EXTENSIONS:
            raise ValueError(f"Unsupported image type {shape.image.ext}")

        # Save image if it doesn't exist
        if not pexists(img_path):
            with open(img_path, "wb") as f:
                f.write(shape.image.blob)

        # Add image style information
        style["img_style"] = {
            "crop_bottom": shape.crop_bottom,
            "crop_top": shape.crop_top,
            "crop_left": shape.crop_left,
            "crop_right": shape.crop_right,
        }

        # Create Picture object
        picture = cls(
            slide_idx,
            shape_idx,
            style,
            [img_path, shape.name, None],  # [img_path, name, caption]
            text_frame,
            slide_area,
            level=level,
        )
        return picture

    def build(self, slide: PPTXSlide) -> PPTXPicture:
        """
        Build the picture in a slide.

        Args:
            slide (PPTXSlide): The slide to build the picture in.

        Returns:
            PPTXPicture: The built picture.
        """
        # Add picture to slide
        shape = slide.shapes.add_picture(
            self.img_path,
            **self.style["shape_bounds"],
        )

        # Set properties
        shape.name = self.style["name"]
        dict_to_object(self.style["img_style"], shape.image)
        apply_fill(shape, self.style["fill_xml"])

        # Apply line style if available
        if self.style["line"] is not None:
            apply_fill(shape.line, self.style["line"]["fill_xml"])
            dict_to_object(self.style["line"], shape.line, exclude=["fill_xml"])

        # Apply shape bounds and rotation
        dict_to_object(self.style["shape_bounds"], shape)
        if hasattr(shape, "rotation"):
            shape.rotation = self.style["rotation"]

        return shape

    @property
    def img_path(self) -> str:
        """
        Get the image path.

        Returns:
            str: The image path.
        """
        return self.data[0]

    @img_path.setter
    def img_path(self, img_path: str) -> None:
        """
        Set the image path.

        Args:
            img_path (str): The image path.
        """
        self.data[0] = img_path

    @property
    def caption(self) -> Optional[str]:
        """
        Get the caption.

        Returns:
            Optional[str]: The caption, or None if not set.
        """
        return self.data[2]

    @caption.setter
    def caption(self, caption: str) -> None:
        """
        Set the caption.

        Args:
            caption (str): The caption.
        """
        self.data[2] = caption

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the picture to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the picture.

        Raises:
            ValueError: If the caption is not found.
        """
        if not style_args.show_image:
            return ""
        if self.caption is None:
            raise ValueError(
                f"Caption not found for picture {self.shape_idx} of slide {self.slide_idx}"
            )
        return (
            self.indent
            + f"<img {self.get_inline_style(style_args)} alt='{self.caption}'/>"
        )


class Placeholder(ShapeElement):
    """
    A class to represent a placeholder shape.
    """

    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: SlidePlaceholder,
        *args,
        **kwargs,
    ) -> Union[Picture, TextBox]:
        """
        Create a Placeholder from a SlidePlaceholder.

        Returns:
            Union[Picture, TextBox]: The created shape element.

        Raises:
            ValueError: If the placeholder type is unsupported.
            AssertionError: If the placeholder has multiple types.
        """
        # Ensure placeholder has only one type
        assert (
            sum(
                [
                    shape.has_text_frame,
                    shape.has_chart,
                    shape.has_table,
                    isinstance(shape, PlaceholderPicture),
                ]
            )
            == 1
        ), "Placeholder should have only one type"

        # Create appropriate shape based on placeholder type
        if isinstance(shape, PlaceholderPicture):
            data = Picture.from_shape(
                slide_idx,
                shape_idx,
                shape,
                *args,
                **kwargs,
            )
        elif shape.has_text_frame:
            data = TextBox.from_shape(
                slide_idx,
                shape_idx,
                shape,
                *args,
                **kwargs,
            )
        else:
            raise ValueError(f"Unsupported placeholder {shape.placeholder_type}")

        return data


class GroupShape(ShapeElement):
    """
    A class to represent a group shape.
    """

    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: PPTXGroupShape,
        style: Dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ) -> "GroupShape":
        """
        Create a GroupShape from a PPTXGroupShape.

        Returns:
            GroupShape: The created GroupShape.
        """
        # Create shape elements for each shape in the group
        data = [
            ShapeElement.from_shape(
                slide_idx,
                (shape_idx + 1) * 100 + i,
                sub_shape,
                config,
                slide_area,
                level=level + 1,
            )
            for i, sub_shape in enumerate(shape.shapes)
        ]

        # Apply shape bounds to each shape in the group
        for idx, shape_bounds in enumerate(parse_groupshape(shape)):
            data[idx].style["shape_bounds"] = shape_bounds

        return cls(
            slide_idx, shape_idx, style, data, text_frame, slide_area, level=level
        )

    def build(self, slide: PPTXSlide) -> PPTXSlide:
        """
        Build the group shape in a slide.

        Args:
            slide (PPTXSlide): The slide to build the group shape in.

        Returns:
            PPTXSlide: The slide with the built group shape.
        """
        for shape in self.data:
            shape.build(slide)
        return slide

    def __iter__(self):
        """
        Iterate over all shapes in the group.

        Yields:
            ShapeElement: Each shape in the group.
        """
        for shape in self.data:
            if isinstance(shape, GroupShape):
                yield from shape
            else:
                yield shape

    def __eq__(self, __value: object) -> bool:
        """
        Check if two group shapes are equal.

        Args:
            __value (object): The object to compare with.

        Returns:
            bool: True if the group shapes are equal, False otherwise.
        """
        if not isinstance(__value, GroupShape) or len(self.data) != len(__value.data):
            return False
        for shape1, shape2 in zip(self.data, __value.data):
            if isinstance(shape1, type(shape2)):
                return False
        return True

    def __repr__(self) -> str:
        """
        Get a string representation of the group shape.

        Returns:
            str: A string representation of the group shape.
        """
        return f"{self.__class__.__name__}: {self.data}"

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the group shape to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the group shape.
        """
        content = "\n".join([shape.to_html(style_args) for shape in self.data])
        if not style_args.show_content:
            content = ""
        return (
            self.indent
            + f"<div {self.get_inline_style(style_args)} data-group-label='{self.group_label}'>\n"
            + content
            + "\n"
            + self.indent
            + "</div>\n"
        )

    @property
    def group_label(self) -> str:
        """
        Get the group label.

        Returns:
            str: The group label.
        """
        return getattr(self, "_group_label", f"group_{self.shape_idx}")

    @group_label.setter
    def group_label(self, value: str) -> None:
        """
        Set the group label.

        Args:
            value (str): The group label.
        """
        self._group_label = value


class FreeShape(ShapeElement):
    """
    A class to represent a free shape.
    """

    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: PPTXAutoShape | PPTXConnector | PPTXGraphicalFrame,
        style: Dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ) -> "FreeShape":
        """
        Create a FreeShape from a PPTXAutoShape.

        Returns:
            FreeShape: The created FreeShape.
        """
        return cls(slide_idx, shape_idx, style, [], text_frame, slide_area, level)

    def to_html(self, style_args: StyleArg) -> str:
        """
        Convert the free shape to HTML.

        Args:
            style_args (StyleArg): The style arguments for HTML conversion.

        Returns:
            str: The HTML representation of the free shape.
        """
        textframe = self.text_frame.to_html(style_args)
        return (
            f"{self.indent}<div {self.get_inline_style(style_args)}>"
            + f"\n{textframe}"
            + f"\n{self.indent}</div>"
        )


class SemanticPicture(ShapeElement):
    """
    A class to represent a semantic picture (table, chart, etc.).
    """

    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: Union[PPTXTable, PPTXChart],
        style: Dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ) -> "SemanticPicture":
        """
        Convert a complex shape to a semantic picture.

        Returns:
            SemanticPicture: The created SemanticPicture.
        """
        shape_type = str(shape.shape_type).split()[0]
        style["img_style"] = {}
        obj = Picture(
            slide_idx,
            shape_idx,
            style,
            [
                "resource/pic_placeholder.png",
                shape.name,
                f"This is a picture of {shape_type}",
            ],
            text_frame,
            slide_area,
            level,
        )
        obj.semantic_name = shape_type
        return obj


# Define shape type mapping
SHAPECAST = {
    MSO_SHAPE_TYPE.AUTO_SHAPE: FreeShape,
    MSO_SHAPE_TYPE.LINE: FreeShape,
    MSO_SHAPE_TYPE.FREEFORM: FreeShape,
    MSO_SHAPE_TYPE.PICTURE: Picture,
    MSO_SHAPE_TYPE.LINKED_PICTURE: Picture,
    MSO_SHAPE_TYPE.PLACEHOLDER: Placeholder,
    MSO_SHAPE_TYPE.GROUP: GroupShape,
    MSO_SHAPE_TYPE.TEXT_BOX: TextBox,
    MSO_SHAPE_TYPE.MEDIA: SemanticPicture,
    MSO_SHAPE_TYPE.TABLE: SemanticPicture,
    MSO_SHAPE_TYPE.CHART: SemanticPicture,
    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT: SemanticPicture,
    MSO_SHAPE_TYPE.LINKED_OLE_OBJECT: SemanticPicture,
    MSO_SHAPE_TYPE.DIAGRAM: SemanticPicture,
    MSO_SHAPE_TYPE.CANVAS: SemanticPicture,
    MSO_SHAPE_TYPE.INK: SemanticPicture,
    MSO_SHAPE_TYPE.IGX_GRAPHIC: SemanticPicture,
    MSO_SHAPE_TYPE.WEB_VIDEO: SemanticPicture,
}
