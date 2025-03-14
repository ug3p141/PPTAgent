import os
import shutil
import subprocess
import tempfile
import traceback
from itertools import product
from time import sleep, time
from types import SimpleNamespace
from typing import Dict, List, Optional, Set, Tuple, Union, Any, Callable

import json_repair
import Levenshtein
from lxml import etree
from pdf2image import convert_from_path
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.text.text import _Paragraph, _Run
from pptx.util import Length, Pt
from tenacity import RetryCallState, retry, stop_after_attempt, wait_fixed

# Set of supported image extensions
IMAGE_EXTENSIONS: Set[str] = {
    "bmp",
    "jpg",
    "jpeg",
    "pgm",
    "png",
    "ppm",
    "tif",
    "tiff",
    "webp",
}

# Common colors and measurements
BLACK = RGBColor(0, 0, 0)
YELLOW = RGBColor(255, 255, 0)
BLUE = RGBColor(0, 0, 255)
BORDER_LEN = Pt(2)
BORDER_OFFSET = Pt(2)
LABEL_LEN = Pt(24)
FONT_LEN = Pt(20)


def is_image_path(file: str) -> bool:
    """
    Check if a file path is an image based on its extension.

    Args:
        file (str): The file path to check.

    Returns:
        bool: True if the file is an image, False otherwise.
    """
    return file.split(".")[-1].lower() in IMAGE_EXTENSIONS


def get_font_style(font: Dict[str, Any]) -> str:
    """
    Convert a font dictionary to a CSS style string.

    Args:
        font (Dict[str, Any]): The font dictionary.

    Returns:
        str: The CSS style string.
    """
    font = SimpleNamespace(**font)
    styles = []

    if hasattr(font, "size") and font.size:
        styles.append(f"font-size: {font.size}pt")

    if hasattr(font, "color") and font.color:
        if all(c in "0123456789abcdefABCDEF" for c in font.color):
            styles.append(f"color: #{font.color}")
        else:
            styles.append(f"color: {font.color}")

    if hasattr(font, "bold") and font.bold:
        styles.append("font-weight: bold")

    if hasattr(font, "italic") and font.italic:
        styles.append("font-style: italic")

    return "; ".join(styles)


def runs_merge(paragraph: _Paragraph) -> Optional[_Run]:
    """
    Merge all runs in a paragraph into a single run.

    Args:
        paragraph (_Paragraph): The paragraph to merge runs in.

    Returns:
        Optional[_Run]: The merged run, or None if there are no runs.
    """
    runs = paragraph.runs

    # Handle field codes
    if len(runs) == 0:
        runs = [
            _Run(r, paragraph)
            for r in parse_xml(paragraph._element.xml.replace("fld", "r")).r_lst
        ]
    if len(runs) == 1:
        return runs[0]
    if len(runs) == 0:
        return None

    # Find the run with the most text
    run = max(runs, key=lambda x: len(x.text))
    run.text = paragraph.text

    # Remove other runs
    for r in runs:
        if r != run:
            r._r.getparent().remove(r._r)
    return run


def older_than(filepath: str, seconds: int = 10, wait: bool = False) -> bool:
    """
    Check if a file is older than a specified number of seconds.

    Args:
        filepath (str): The path to the file.
        seconds (int): The number of seconds to check against.
        wait (bool): Whether to wait for the file to exist.

    Returns:
        bool: True if the file is older than the specified number of seconds, False otherwise.
    """
    if not os.path.exists(filepath):
        while wait:
            print("waiting for:", filepath)
            sleep(1)
            if os.path.exists(filepath):
                sleep(seconds)
                return True
        return False
    file_creation_time = os.path.getctime(filepath)
    current_time = time()
    return seconds < (current_time - file_creation_time)


def edit_distance(text1: str, text2: str) -> float:
    """
    Calculate the normalized edit distance between two strings.

    Args:
        text1 (str): The first string.
        text2 (str): The second string.

    Returns:
        float: The normalized edit distance (0.0 to 1.0, where 1.0 means identical).
    """
    if not text1 and not text2:
        return 1.0
    return 1 - Levenshtein.distance(text1, text2) / max(len(text1), len(text2))


def get_slide_content(
    doc_json: Dict[str, Any], slide_title: str, slide: Dict[str, Any]
) -> str:
    """
    Get the content for a slide based on its title and description.

    Args:
        doc_json (Dict[str, Any]): The document JSON.
        slide_title (str): The title of the slide.
        slide (Dict[str, Any]): The slide data.

    Returns:
        str: The slide content.
    """
    slide_desc = slide.get("description", "")
    slide_content = f"Slide Purpose: {slide_title}\nSlide Description: {slide_desc}\n"

    for key in slide.get("subsections", []):
        slide_content += "Slide Content Source: "

        for section in doc_json["sections"]:
            subsections = section.get("subsections", [])

            # Handle dictionary subsections
            if isinstance(subsections, dict) and len(subsections) == 1:
                subsections = [
                    {"title": k, "content": v} for k, v in subsections.items()
                ]

            for subsection in subsections:
                try:
                    if edit_distance(key, subsection["title"]) > 0.8:
                        slide_content += f"# {key} \n{subsection['content']}\n"
                except Exception as e:
                    print(f"Error processing subsection: {e}")

    return slide_content


def tenacity_log(retry_state: RetryCallState) -> None:
    """
    Log function for tenacity retries.

    Args:
        retry_state (RetryCallState): The retry state.
    """
    print(retry_state)
    traceback.print_tb(retry_state.outcome.exception().__traceback__)


def get_json_from_response(response: str) -> Dict[str, Any]:
    """
    Extract JSON from a text response.

    Args:
        response (str): The response text.

    Returns:
        Dict[str, Any]: The extracted JSON.

    Raises:
        Exception: If JSON cannot be extracted from the response.
    """
    response = response.strip()

    # Try to extract JSON from markdown code blocks
    l, r = response.rfind("```json"), response.rfind("```")
    if l != -1 and r != -1:
        json_obj = json_repair.loads(response[l + 7 : r].strip())
        if isinstance(json_obj, (dict, list)):
            return json_obj

    # Try to find JSON by looking for matching braces
    open_braces = []
    close_braces = []

    for i, char in enumerate(response):
        if char == "{":
            open_braces.append(i)
        elif char == "}":
            close_braces.append(i)

    for i, j in product(open_braces, reversed(close_braces)):
        if i > j:
            continue
        try:
            json_obj = json_repair.loads(response[i : j + 1])
            if isinstance(json_obj, (dict, list)):
                return json_obj
        except Exception:
            pass

    raise Exception("JSON not found in the given output", response)


# Create a tenacity decorator with custom settings
tenacity = retry(
    wait=wait_fixed(3), stop=stop_after_attempt(5), after=tenacity_log, reraise=True
)


@tenacity
def ppt_to_images(pptx: str, output_dir: str) -> None:
    """
    Convert a PowerPoint file to images.

    Args:
        pptx (str): The path to the PowerPoint file.
        output_dir (str): The directory to save the images to.

    Raises:
        AssertionError: If the file does not exist or conversion fails.
    """
    assert pexists(pptx), f"File {pptx} does not exist"
    os.makedirs(output_dir, exist_ok=True)

    with tempfile.NamedTemporaryFile(suffix=".pdf") as temp_pdf:
        # Convert PPTX to PDF
        command_list = [
            "unoconvert",
            "--convert-to",
            "pdf",
            pptx,
            temp_pdf.name,
        ]

        try:
            subprocess.run(command_list, check=True, stdout=subprocess.DEVNULL)
        except subprocess.CalledProcessError as e:
            raise AssertionError(f"PPTX to PDF conversion failed: {e}")

        assert pexists(
            temp_pdf.name
        ), "PPTX convert failed, check the installation of unoserver"

        # Convert PDF to images
        try:
            images = convert_from_path(temp_pdf.name, dpi=72)
            for i, img in enumerate(images):
                img.save(pjoin(output_dir, f"slide_{i+1:04d}.jpg"))
        except Exception as e:
            raise AssertionError(f"PDF to image conversion failed: {e}")


@tenacity
def wmf_to_images(blob: bytes, filepath: str) -> None:
    """
    Convert a WMF blob to an image.

    Args:
        blob (bytes): The WMF blob.
        filepath (str): The path to save the image to.

    Raises:
        AssertionError: If the conversion fails.
    """
    with tempfile.NamedTemporaryFile(suffix=".wmf") as temp_wmf:
        # Write blob to temporary file
        with open(temp_wmf.name, "wb") as f:
            f.write(blob)

        # Convert WMF to JPG
        command_list = [
            "unoconvert",
            "--convert-to",
            "jpg",
            temp_wmf.name,
            filepath,
        ]

        try:
            subprocess.run(command_list, check=True, stdout=subprocess.DEVNULL)
        except subprocess.CalledProcessError as e:
            raise AssertionError(f"WMF conversion failed: {e}")

    assert pexists(filepath), f"WMF convert failed, output file {filepath} not found"


def extract_fill(shape: BaseShape) -> Optional[Tuple[str, str]]:
    """
    Extract fill information from a shape.

    Args:
        shape (BaseShape): The shape to extract fill from.

    Returns:
        Optional[Tuple[str, str]]: The fill string and XML, or None if the shape has no fill.
    """
    if not hasattr(shape, "fill"):
        return None

    try:
        fill_str = "Fill: " + str(shape.fill.value)
        fill_xml = shape.fill._xPr.xml
        return fill_str, fill_xml
    except Exception as e:
        print(f"Error extracting fill: {e}")
        return None


def apply_fill(shape: BaseShape, fill_xml: Optional[str]) -> None:
    """
    Apply fill XML to a shape.

    Args:
        shape (BaseShape): The shape to apply fill to.
        fill_xml (Optional[str]): The fill XML to apply.
    """
    if fill_xml is None:
        return

    try:
        new_element = etree.fromstring(fill_xml)
        shape.fill._xPr.getparent().replace(shape.fill._xPr, new_element)
    except Exception as e:
        print(f"Error applying fill: {e}")


def parse_groupshape(groupshape: GroupShape) -> List[Dict[str, Length]]:
    """
    Parse a group shape to get the bounds of its child shapes.

    Args:
        groupshape (GroupShape): The group shape to parse.

    Returns:
        List[Dict[str, Length]]: The bounds of the child shapes.

    Raises:
        AssertionError: If the input is not a GroupShape.
    """
    assert isinstance(groupshape, GroupShape), "Input must be a GroupShape"

    # Get group bounds
    group_top_left_x = groupshape.left
    group_top_left_y = groupshape.top
    group_width = groupshape.width
    group_height = groupshape.height

    # Get shape bounds
    shape_top_left_x = min([sp.left for sp in groupshape.shapes])
    shape_top_left_y = min([sp.top for sp in groupshape.shapes])
    shape_width = (
        max([sp.left + sp.width for sp in groupshape.shapes]) - shape_top_left_x
    )
    shape_height = (
        max([sp.top + sp.height for sp in groupshape.shapes]) - shape_top_left_y
    )

    # Calculate bounds for each shape in the group
    group_shape_xy = []
    for sp in groupshape.shapes:
        group_shape_left = (
            sp.left - shape_top_left_x
        ) * group_width / shape_width + group_top_left_x
        group_shape_top = (
            sp.top - shape_top_left_y
        ) * group_height / shape_height + group_top_left_y
        group_shape_width = sp.width * group_width / shape_width
        group_shape_height = sp.height * group_height / shape_height

        group_shape_xy.append(
            {
                "left": Length(group_shape_left),
                "top": Length(group_shape_top),
                "width": Length(group_shape_width),
                "height": Length(group_shape_height),
            }
        )

    return group_shape_xy


def is_primitive(obj: Any) -> bool:
    """
    Check if an object is a primitive type or a collection of primitive types.

    Args:
        obj (Any): The object to check.

    Returns:
        bool: True if the object is a primitive type or a collection of primitive types, False otherwise.
    """
    if isinstance(obj, (list, tuple, set, frozenset)):
        return all(is_primitive(item) for item in obj)

    return isinstance(
        obj, (int, float, complex, bool, str, bytes, bytearray, type(None))
    )


DEFAULT_EXCLUDE: Set[str] = set(["element", "language_id", "ln", "placeholder_format"])


def object_to_dict(
    obj: Any,
    result: Optional[Dict[str, Any]] = None,
    exclude: Optional[Set[str]] = None,
) -> Dict[str, Any]:
    """
    Convert an object to a dictionary.

    Args:
        obj (Any): The object to convert.
        result (Optional[Dict[str, Any]]): The result dictionary to update.
        exclude (Optional[Set[str]]): The attributes to exclude.

    Returns:
        Dict[str, Any]: The dictionary representation of the object.
    """
    if result is None:
        result = {}

    exclude = DEFAULT_EXCLUDE.union(exclude or set())

    for attr in dir(obj):
        if attr in exclude or attr.startswith("_") or callable(getattr(obj, attr)):
            continue

        try:
            attr_value = getattr(obj, attr)

            # Handle complex numbers
            if hasattr(attr_value, "real"):
                attr_value = attr_value.real

            # Convert size to points
            if attr == "size" and isinstance(attr_value, int):
                attr_value = Length(attr_value).pt

            # Only include primitive types
            if is_primitive(attr_value):
                result[attr] = attr_value
        except Exception:
            pass

    return result


def merge_dict(d1: Dict[str, Any], d2: List[Dict[str, Any]]) -> Dict[str, Any]:
    """
    Merge multiple dictionaries, keeping only values that are the same across all dictionaries.

    Args:
        d1 (Dict[str, Any]): The base dictionary.
        d2 (List[Dict[str, Any]]): The dictionaries to merge.

    Returns:
        Dict[str, Any]: The merged dictionary.
    """
    if not d2:
        return d1

    for key in list(d1.keys()):
        values = [d[key] for d in d2]
        if d1[key] is not None and len(values) != 1:
            values.append(d1[key])

        # Skip if any value is None or values are not all the same
        if (
            not values
            or values[0] is None
            or not all(value == values[0] for value in values)
        ):
            continue

        # Set the value in the base dictionary and clear it in the other dictionaries
        d1[key] = values[0]
        for d in d2:
            d[key] = None
    return d1


def dict_to_object(
    dict_obj: Dict[str, Any], obj: Any, exclude: Optional[Set[str]] = None
) -> None:
    """
    Apply dictionary values to an object.

    Args:
        dict_obj (Dict[str, Any]): The dictionary with values to apply.
        obj (Any): The object to apply values to.
        exclude (Optional[Set[str]]): The keys to exclude.
    """
    if exclude is None:
        exclude = set()

    for key, value in dict_obj.items():
        if key not in exclude and value is not None:
            try:
                setattr(obj, key, value)
            except Exception as e:
                print(f"Error setting attribute {key}: {e}")


class Config:
    """
    Configuration class for the application.
    """

    def __init__(
        self,
        rundir: Optional[str] = None,
        session_id: Optional[str] = None,
        debug: bool = True,
    ):
        """
        Initialize the configuration.

        Args:
            rundir (Optional[str]): The run directory.
            session_id (Optional[str]): The session ID.
            debug (bool): Whether to enable debug mode.
        """
        self.DEBUG = debug

        if session_id is not None:
            self.set_session(session_id)
        if rundir is not None:
            self.set_rundir(rundir)

    def set_session(self, session_id: str) -> None:
        """
        Set the session ID and update the run directory.

        Args:
            session_id (str): The session ID.
        """
        self.session_id = session_id
        self.set_rundir(f"./runs/{session_id}")

    def set_rundir(self, rundir: str) -> None:
        """
        Set the run directory and create necessary subdirectories.

        Args:
            rundir (str): The run directory.
        """
        self.RUN_DIR = rundir
        self.IMAGE_DIR = pjoin(self.RUN_DIR, "images")

        for the_dir in [self.RUN_DIR, self.IMAGE_DIR]:
            os.makedirs(the_dir, exist_ok=True)

    def set_debug(self, debug: bool) -> None:
        """
        Set the debug mode.

        Args:
            debug (bool): Whether to enable debug mode.
        """
        self.DEBUG = debug

    def remove_rundir(self) -> None:
        """
        Remove the run directory and its subdirectories.
        """
        if pexists(self.RUN_DIR):
            shutil.rmtree(self.RUN_DIR)
        if pexists(self.IMAGE_DIR):
            shutil.rmtree(self.IMAGE_DIR)

    def __repr__(self) -> str:
        """
        Get a string representation of the configuration.

        Returns:
            str: A string representation of the configuration.
        """
        attrs = []
        for attr in dir(self):
            if not attr.startswith("_") and not callable(getattr(self, attr)):
                attrs.append(f"{attr}={getattr(self, attr)}")
        return f"Config({', '.join(attrs)})"


# Path utility functions
pjoin = os.path.join
pexists = os.path.exists
pbasename = os.path.basename
pdirname = os.path.dirname
