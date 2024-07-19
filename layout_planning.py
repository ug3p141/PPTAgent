from types import SimpleNamespace
from typing import List
from matplotlib import patches, pyplot as plt
from presentation import TextFrame, Presentation, ShapeElement
from pptx.text.fonts import FontFiles
from PIL.ImageFont import truetype
from utils import pjoin, base_config


# 根据已经生成好的template
class LayoutPlanner:
    def __init__(self, shapes: List[ShapeElement], width: int, height: int):
        self.font_files = FontFiles._installed_fonts()
        self.shapes = shapes
        self.width = width
        self.height = height

    def get_html(self):
        pass

    def plan(self, scene):
        pass

    # 如何解决font family是None的问题
    def calc_textframe(self, textframe: TextFrame):
        width, height = textframe.content_width, textframe.content_height
        for para in textframe.data:
            font = SimpleNamespace(**para["font"])
            # 自定义一个类然后使用bisect 进行比较和查询
            font = truetype(self.find_font(font), point_size)

    def draw(self):
        fig, ax = plt.subplots()

        for shape in self.shapes:
            rect = patches.Rectangle(
                (shape.left, shape.top),
                shape.width,
                shape.height,
                linewidth=1,
                edgecolor="r",
                facecolor="none",
            )
            ax.add_patch(rect)

        ax.set_xlim(0, self.width)
        ax.set_ylim(0, self.height)
        ax.set_aspect("equal")

        plt.gca().invert_yaxis()  # 将Y轴翻转，以匹配HTML中的坐标系
        plt.show()

    def find_font(self, font: SimpleNamespace):
        pass


if __name__ == "__main__":
    prs = Presentation.from_file(
        pjoin(
            base_config.PPT_DIR,
            "中文信息联合党支部2022年述职报告.pptx",
        )
    )

    # 创建 BoxModel 对象并绘制图像
    box_model = LayoutPlanner(
        [i for i in prs.slides[0].shapes if not i.is_background],
        prs.slide_width,
        prs.slide_height,
    )
    box_model.draw()

""" 三种计算方式：
import freetype

def get_text_dimensions(text, font_path, font_size, line_spacing=1.0, is_bold=False, is_italic=False):
    face = freetype.Face(font_path)
    face.set_char_size(font_size * 64)

    # Optional: Set style attributes
    if is_bold:
        face.style_flags |= freetype.FT_STYLE_FLAG_BOLD
    if is_italic:
        face.style_flags |= freetype.FT_STYLE_FLAG_ITALIC

    width = 0
    max_height = 0
    max_line_height = 0
    pen_x = 0
    pen_y = 0

    for char in text:
        face.load_char(char)
        bitmap = face.glyph.bitmap

        pen_x += face.glyph.advance.x // 64
        max_line_height = max(max_line_height, face.glyph.metrics.height // 64)
        max_height = max(max_height, bitmap.rows)

        if char == '\n':
            width = max(width, pen_x)
            pen_x = 0
            pen_y += max_line_height * line_spacing
            max_line_height = 0

    width = max(width, pen_x)
    height = pen_y + max_line_height

    return width, height

# Usage
font_path = "/path/to/your/font.ttf"
font_size = 12 * 72 // 96  # Points to pixels conversion
text = "Hello, World!\nThis is a test."

width, height = get_text_dimensions(text, font_path, font_size, line_spacing=1.2, is_bold=True, is_italic=False)
print(f"Text width: {width}, Text height: {height}")



from PIL import ImageFont, ImageDraw, Image

def get_text_dimensions(text, font_path, font_size, line_spacing=1.0):
    font = ImageFont.truetype(font_path, font_size)

    # Create a dummy image to get a drawing context
    dummy_image = Image.new("RGB", (1, 1))
    draw = ImageDraw.Draw(dummy_image)

    lines = text.split('\n')
    width = max(draw.textsize(line, font=font)[0] for line in lines)
    height = sum(draw.textsize(line, font=font)[1] for line in lines) + (len(lines) - 1) * int(font_size * line_spacing)

    return width, height

# Usage
font_path = "/path/to/your/font.ttf"
font_size = 12
text = "Hello, World!\nThis is a test."

width, height = get_text_dimensions(text, font_path, font_size, line_spacing=1.2)
print(f"Text width: {width}, Text height: {height}")



from PyQt5.QtWidgets import QApplication
from PyQt5.QtGui import QFont, QFontMetrics

def get_text_dimensions(text, font_family, font_size, is_bold=False, is_italic=False, line_spacing=1.0):
    app = QApplication([])

    font = QFont(font_family, font_size)
    font.setBold(is_bold)
    font.setItalic(is_italic)

    metrics = QFontMetrics(font)
    lines = text.split('\n')
    width = max(metrics.width(line) for line in lines)
    height = sum(metrics.height() for _ in lines) + (len(lines) - 1) * metrics.lineSpacing() * (line_spacing - 1)

    return width, height

# Usage
font_family = "Arial"
font_size = 12
text = "Hello, World!\nThis is a test."

width, height = get_text_dimensions(text, font_family, font_size, is_bold=True, is_italic=False, line_spacing=1.2)
print(f"Text width: {width}, Text height: {height}")




"""
