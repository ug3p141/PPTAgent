from itertools import count
from typing import Iterator

from pptx import Presentation as PPTXPre
from pptx.chart.chart import Chart as PPTXChart
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
from pptx.shapes.autoshape import Shape as PPTXAutoShape
from pptx.shapes.base import BaseShape
from pptx.shapes.connector import Connector as PPTXConnector
from pptx.shapes.group import GroupShape as PPTXGroupShape
from pptx.shapes.picture import Picture as PPTXPicture
from pptx.shapes.placeholder import PlaceholderPicture, SlidePlaceholder
from pptx.slide import Slide as PPTXSlide
from pptx.table import Table as PPTXTable
from pptx.text.text import _Paragraph, _Run
from pptx.util import Pt
from rich import print

from utils import (
    IMAGE_EXTENSIONS,
    Config,
    apply_fill,
    dict_to_object,
    extract_fill,
    get_font_style,
    merge_dict,
    object_to_dict,
    parse_groupshape,
    pexists,
    pjoin,
    runs_merge,
)

# slide, group, shape, textframe, ul, run
INDENT = "\t"


# textframe: shape bounds
# paragraph: space, alignment, level, font bullet #  这里都不允许改动
# run: font, hyperlink, text
class Run:
    def __init__(self, run: _Run, counter: Iterator, para_font: dict):
        self.idx = next(counter)
        self.text = run.text
        self.font = {}
        self.font = merge_dict(object_to_dict(run.font), para_font)
        if self.font["size"] is None:
            self.font["size"] = Pt(18)  # default size

    # ? 这里似乎不需要style因为后面会css化?
    def to_html(self, style_kwargs, tag: str):
        id_str = f" id='{self.idx}'" if style_kwargs.get("textframe_id", False) else ""
        style = (
            f" style='{get_font_style(self.font)}'"
            if style_kwargs.get("font_style", False)
            else ""
        )
        return f"<{tag}{id_str}{style}>{self.text}</{tag}>"


class Paragraph:
    def __init__(self, paragraph: _Paragraph, counter: Iterator, textframe_font: dict):
        runs = runs_merge(paragraph)
        self.is_paragraph = len(runs) > 0
        para_font = merge_dict(object_to_dict(paragraph.font), textframe_font)
        self.runs = [Run(r, counter, para_font) for r in runs]
        self.bullet = paragraph.bullet

    def to_html(self, style_kwargs):
        repr_list = []
        run_tag = "li" if self.bullet else "p"
        for run in self.runs:
            repr_list.append(run.to_html(style_kwargs, run_tag))
        return repr_list


class TextFrame:
    def __init__(self, shape: BaseShape, level: int):
        if not shape.has_text_frame:
            self.is_textframe = False
            return
        self.is_textframe = True
        self.level = level
        self.text = shape.text
        textframe_font = object_to_dict(shape.text_frame.font)
        counter = count(1)
        paragraphs = [
            Paragraph(paragraph, counter, textframe_font)
            for paragraph in shape.text_frame.paragraphs
        ]
        self.paragraphs = [para for para in paragraphs if para.is_paragraph]

    def to_html(self, style_kwargs):
        if not self.is_textframe:
            return ""
        repr_list = []
        pre_bullet = None
        bullets = []
        for para in self.paragraphs:
            if (para.bullet != pre_bullet or para == self.paragraphs[-1]) and len(
                bullets
            ) != 0:
                repr_list += ["<ul>"] + [INDENT + repr for repr in bullets] + ["</ul>"]
            if para.bullet is None:
                repr_list.extend(para.to_html(style_kwargs))
            else:
                bullets.extend(para.to_html(style_kwargs))
            pre_bullet = para.bullet
        return "\n".join([INDENT * self.level + repr for repr in repr_list])

    def __repr__(self):
        if not self.is_textframe:
            return "TextFrame: null"
        return f"TextFrame: {self.text}"

    def __len__(self):
        if not self.is_textframe:
            return 0
        return len(self.text)


class ShapeElement:
    def __init__(
        self,
        slide_idx: int,
        shape_idx: int,
        style: dict,
        data: dict,
        text_frame: TextFrame,
        slide_area: float,
        level: int,
    ):
        self.slide_idx = slide_idx
        self.shape_idx = shape_idx
        self.style = style
        self.data = data
        self.text_frame = text_frame
        self.closures = {}
        self.slide_area = slide_area
        self.level = level

    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: BaseShape,
        config: Config,
        slide_area: float,
        level: int = 0,
    ):
        if shape_idx > 100 and isinstance(shape, PPTXGroupShape):
            raise ValueError(f"nested group shapes are not allowed")
        line = None
        if "line" in dir(shape) and shape.line._ln is not None:
            line = {
                "fill": extract_fill(shape.line),
                "width": shape.line.width,
                "dash_style": shape.line.dash_style,
            }
        fill = extract_fill(shape)
        style = {
            "shape_bounds": {
                "width": shape.width,
                "height": shape.height,
                "left": shape.left,
                "top": shape.top,
            },
            "shape_type": str(shape.shape_type).split("(")[0].lower(),
            "rotation": shape.rotation,
            "fill": fill,
            "line": line,
        }
        text_frame = TextFrame(shape, level + 1)
        obj = SHAPECAST.get(shape.shape_type, UnsupportedShape).from_shape(
            slide_idx,
            shape_idx,
            shape,
            style,
            text_frame,
            config,
            slide_area,
            level,
        )
        obj.xml = shape._element.xml
        # ? obj.shape = shape
        return obj

    def build(self, slide, shape):
        apply_fill(shape, self.style["fill"])
        if self.style["line"] is not None:
            apply_fill(shape.line, self.style["line"]["fill"])
            dict_to_object(self.style["line"], shape.line, exclude=["fill"])

        dict_to_object(self.style["shape_bounds"], shape)
        if "rotation" in dir(shape):
            shape.rotation = self.style["rotation"]
        return shape

    def __repr__(self) -> str:
        return f"{self.__class__.__name__}: shape {self.shape_idx} of slide {self.slide_idx}"

    def to_html(self, **kwargs) -> str:
        return ""

    @property
    def indent(self):
        return "\t" * self.level

    @property
    def left(self):
        return self.style["shape_bounds"]["left"].pt

    @left.setter
    def left(self, value):
        self.style["shape_bounds"]["left"] = value

    @property
    def top(self):
        return self.style["shape_bounds"]["top"].pt

    @top.setter
    def top(self, value):
        self.style["shape_bounds"]["top"] = value

    @property
    def width(self):
        return self.style["shape_bounds"]["width"].pt

    @width.setter
    def width(self, value):
        self.style["shape_bounds"]["width"] = value

    @property
    def height(self):
        return self.style["shape_bounds"]["height"].pt

    @height.setter
    def height(self, value):
        self.style["shape_bounds"]["height"] = value

    @property
    def area(self):
        return self.width * self.height

    def get_inline_style(self, style_kwargs: dict):
        id_str = (
            f" id='{self.shape_idx}'" if style_kwargs.get("element_id", False) else ""
        )
        styles = []
        if style_kwargs.get("area", False):
            styles.append(f" data-relative-area='{self.area*100/self.slide_area:.2f}%'")
        if style_kwargs.get("geometry", False):
            styles.append(f"left: {self.left}pt; top: {self.top}pt;")
        if style_kwargs.get("size", False):
            styles.append(f"width: {self.width}pt; height: {self.height}pt;")

        if len(styles) != 0:
            return id_str + " style='" + " ".join(styles) + "'"
        return id_str


class UnsupportedShape(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: BaseShape,
        *args,
        **kwargs,
    ):
        raise ValueError(f"unsupported shape {shape.shape_type}")


class TextBox(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: TextFrame,
        style: dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ):
        return cls(slide_idx, shape_idx, style, None, text_frame, slide_area, level)

    def to_html(self, **style_kwargs) -> str:

        return (
            f"{self.indent}<div{self.get_inline_style(style_kwargs)}>\n"
            + self.text_frame.to_html(style_kwargs)
            + f"\n{self.indent}</div>\n"
        )


class Picture(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: PPTXPicture,
        style: dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ):
        if shape.image.ext not in IMAGE_EXTENSIONS:
            raise ValueError(f"unsupported image type {shape.image.ext}")
        img_path = pjoin(
            config.IMAGE_DIR,
            f"{shape.image.sha1}.{shape.image.ext}",
        )
        if not pexists(img_path):
            with open(img_path, "wb") as f:
                f.write(shape.image.blob)
        style["img_style"] = {
            "crop_bottom": shape.crop_bottom,
            "crop_top": shape.crop_top,
            "crop_left": shape.crop_left,
            "crop_right": shape.crop_right,
        }
        picture = cls(
            slide_idx,
            shape_idx,
            style,
            [img_path, shape.name, ""],
            text_frame,
            slide_area,
            level=level,
        )
        return picture

    def build(self, slide: PPTXSlide):
        shape = slide.shapes.add_picture(
            self.img_path,
            **self.style["shape_bounds"],
        )
        shape.name = self.data[1]
        dict_to_object(self.style["img_style"], shape.image)
        return super().build(slide, shape)

    @property
    def img_path(self):
        return self.data[0]

    @img_path.setter
    def img_path(self, img_path: str):
        self.data[0] = img_path

    @property
    def caption(self):
        return self.data[2]

    @caption.setter
    def caption(self, caption: str):
        self.data[2] = caption

    def to_html(self, **style_kwargs) -> str:
        if style_kwargs.get("image_id", False):
            style_kwargs["element_id"] = True
        if not self.caption:
            return ""
        return (
            self.indent
            + f"<figure{self.get_inline_style(style_kwargs)}>\n"
            + f"{self.indent+INDENT}<figcaption>{self.caption}</figcaption>\n"
            + self.indent
            + "</figure>\n"
        )


class Placeholder(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: SlidePlaceholder,
        style: dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ):
        assert (
            sum(
                [
                    shape.has_text_frame,
                    shape.has_chart,
                    shape.has_table,
                    isinstance(shape, PlaceholderPicture),
                ]
            )
            == 1
        ), "placeholder should have only one type"
        if isinstance(shape, PlaceholderPicture):
            data = Picture.from_shape(
                slide_idx,
                shape_idx,
                shape,
                style,
                text_frame,
                config,
                slide_area,
                level,
            )
        elif shape.has_text_frame:
            data = TextBox.from_shape(
                slide_idx,
                shape_idx,
                shape,
                style,
                text_frame,
                config,
                slide_area,
                level,
            )
        elif shape.has_chart or shape.has_table:
            data = GraphicalShape.from_shape(
                slide_idx,
                shape_idx,
                shape,
                style,
                text_frame,
                config,
                slide_area,
                level,
            )
        return data


class GroupShape(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: PPTXGroupShape,
        style: dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ):
        data = [
            ShapeElement.from_shape(
                slide_idx,
                (shape_idx + 1) * 100 + i,
                sub_shape,
                config,
                slide_area,
                level=level + 1,
            )
            for i, sub_shape in enumerate(shape.shapes)
        ]
        for idx, shape_bounds in enumerate(parse_groupshape(shape)):
            data[idx].style["shape_bounds"] = shape_bounds
        return cls(
            slide_idx, shape_idx, style, data, text_frame, slide_area, level=level
        )

    def build(self, slide: PPTXSlide):
        for sub_shape in self.data:
            if isinstance(sub_shape, (Picture, GroupShape)):
                new_shape = sub_shape.build(slide)
            else:
                new_shape = slide.shapes._spTree.insert_element_before(
                    parse_xml(sub_shape.xml), "p:extLst"
                )
            for closure in sub_shape.closures.values():
                closure(new_shape)

    def __iter__(self):
        for shape in self.data:
            if isinstance(shape, GroupShape):
                yield from shape
            else:
                yield shape

    def __eq__(self, __value: object) -> bool:
        if not isinstance(__value, GroupShape) or len(self.data) != len(__value.data):
            return False
        for shape1, shape2 in zip(self.data, __value.data):
            if isinstance(shape1, type(shape2)):
                return False
        return True

    def __repr__(self) -> str:
        return f"{self.__class__.__name__}: {self.data}"

    def to_html(self, **style_kwargs) -> str:
        return (
            self.indent
            + f"<div class='{self.group_label}'{self.get_inline_style(style_kwargs)}>\n"
            + "\n".join([shape.to_html(**style_kwargs) for shape in self.data])
            + "\n"
            + self.indent
            + "</div>\n"
        )


class GraphicalShape(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: PPTXChart | PPTXTable,
        style: dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ):
        return cls(
            slide_idx, shape_idx, style, None, text_frame, slide_area, level=level
        )

    @property
    def orig_shape(self):
        return self.style["shape_type"].strip()

    def normalize(self):
        shape = Picture(
            self.slide_idx,
            self.shape_idx,
            {"img_style": {}} | self.style,
            [
                "resource/pic_placeholder.png",
                self.orig_shape,
                f"{self.orig_shape}_{self.shape_idx}",
            ],
            self.text_frame,
            self.slide_area,
            level=self.level,
        )
        shape.style["fill"] = None
        shape.style["line"] = None
        shape.text_frame.is_textframe = False
        return shape


class FreeShape(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: PPTXAutoShape,
        style: dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ):
        data = {
            "shape_type": shape.auto_shape_type.real,
            "svg_tag": str(shape.auto_shape_type).split()[0].lower(),
        }
        return cls(
            slide_idx, shape_idx, style, data, text_frame, slide_area, level=level
        )

    def to_html(self, **style_kwargs) -> str:
        textframe = self.text_frame.to_html(style_kwargs)
        return (
            f"{self.indent}<div data-shape-type='{self.data['svg_tag']}'{self.get_inline_style(style_kwargs)}>"
            + f"\n{textframe}"
            + f"\n{self.indent}</div>"
        )


class Connector(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: PPTXConnector,
        style: dict,
        text_frame: TextFrame,
        config: Config,
        slide_area: float,
        level: int,
    ):
        return FreeShape(
            slide_idx,
            shape_idx,
            style,
            {"shape_type": "connector", "svg_tag": "connector"},
            text_frame,
            slide_area,
            level,
        )


class SlidePage:
    def __init__(
        self,
        shapes: list[ShapeElement],
        slide_idx: int,
        real_idx: int,
        slide_notes: str,
        slide_layout_name: str,
        slide_title: str,
        slide_width: int,
        slide_height: int,
    ):
        self.shapes = shapes
        self.slide_idx = slide_idx
        self.real_idx = real_idx
        self.slide_notes = slide_notes
        self.slide_layout_name = slide_layout_name
        self.slide_title = slide_title
        self.slide_width = slide_width
        self.slide_height = slide_height
        groups_shapes_labels = []
        for shape in self.shape_filter(GroupShape):
            for group_shape in groups_shapes_labels:
                if group_shape == shape:
                    shape.group_label = group_shape.group_label
                    continue
                groups_shapes_labels.append(shape)
                shape.group_label = f"group_{len(groups_shapes_labels)}"

    @classmethod
    def from_slide(
        cls,
        slide: PPTXSlide,
        slide_idx: int,
        real_idx: int,
        slide_width: int,
        slide_height: int,
        config: Config,
    ):
        shapes = [
            ShapeElement.from_shape(
                slide_idx, i, shape, config, slide_width * slide_height
            )
            for i, shape in enumerate(slide.shapes)
            if shape.visible
        ]
        slide_layout_name = slide.slide_layout.name if slide.slide_layout else None
        slide_title = slide.shapes.title.text if slide.shapes.title else None
        slide_notes = (
            slide.notes_slide.notes_text_frame.text
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame
            else None
        )
        return cls(
            shapes,
            slide_idx,
            real_idx,
            slide_notes,
            slide_layout_name,
            slide_title,
            slide_width,
            slide_height,
        )

    # TODO 在build的时候可以根据class 建议一个 class, shapes mapping
    def build(self, slide: PPTXSlide):
        for ph in slide.placeholders:
            ph.element.getparent().remove(ph.element)

        for shape in self.shapes:
            if isinstance(shape, (Picture, GroupShape)):
                shape.build(slide)
            else:
                slide.shapes._spTree.insert_element_before(
                    parse_xml(shape.xml), "p:extLst"
                )
            for run in shape.text_frame:
                run.get("closure", lambda x: x)(slide.shapes[-1])
            shape.get("closure", lambda x: x)(slide.shapes[-1])
        return slide

    def shape_filter(self, shape_type: type, shapes: list[ShapeElement] = None):
        if shapes is None:
            shapes = self.shapes
        for shape in shapes:
            if isinstance(shape, shape_type):
                yield shape
            elif isinstance(shape, GroupShape):
                yield from self.shape_filter(shape_type, shape.data)

    def get_content_types(self, shapes: list[ShapeElement] = None):
        content_types = set()
        if shapes is None:
            shapes = self.shapes
        for shape in shapes:
            if isinstance(shape, GraphicalShape):
                content_types.add(shape.orig_shape)
            elif isinstance(shape, Picture):
                if not shape.img_path == "resource/pic_placeholder.png":
                    content_types.add("picture")
                else:

                    content_types.add(shape.caption)
            elif isinstance(shape, GroupShape):
                content_types.union(self.get_content_types(shape.data))
        return sorted(list(content_types))

    def to_html(self, **kwargs) -> str:
        return "".join(
            [
                "<!DOCTYPE html>\n<html>\n",
                (f"<title>{self.slide_title}</title>\n" if self.slide_title else ""),
                f'<body style="width:{self.slide_width}pt; height:{self.slide_height}pt;">\n',
                "\n".join([shape.to_html(**kwargs) for shape in self.shapes]),
                "</body>\n</html>\n",
            ]
        )

    # TODO 将shapes 转换成一个class , data:shapes, requirements 的list
    def to_tree(self, template: dict):
        pass

    def to_text(self, show_image: bool = True) -> str:
        return "\n".join(
            [
                shape.text_frame.text.strip()
                for shape in self.shapes
                if shape.text_frame.is_textframe
            ]
            + [
                "Image: " + shape.caption
                for shape in self.shape_filter(Picture)
                if show_image
            ]
        )

    def normalize(self, shapes: list[ShapeElement] = None):
        if shapes is None:
            shapes = self.shapes
        for shape_idx, shape in enumerate(shapes):
            if isinstance(shape, GraphicalShape):
                shapes[shape_idx] = shape.normalize()
            elif isinstance(shape, GroupShape):
                self.normalize(shape.data)

    @property
    def text_length(self):
        return sum([len(shape.text_frame) for shape in self.shapes])

    def __iter__(self):
        for shape in self.shapes:
            if isinstance(shape, GroupShape):
                yield from shape
            else:
                yield shape

    def __len__(self):
        return len(self.shapes)


class Presentation:
    def __init__(
        self,
        slides: list[SlidePage],
        error_history: list[str],
        slide_width: float,
        slide_height: float,
        file_path: str,
        num_pages: int,
    ) -> None:
        self.slides = slides
        self.error_history = error_history
        self.slide_width = slide_width
        self.slide_height = slide_height
        self.num_pages = num_pages
        self.source_file = file_path
        self.prs = PPTXPre(self.source_file)
        self.layout_mapping = {layout.name: layout for layout in self.prs.slide_layouts}
        self.prs.core_properties.last_modified_by = "PPTAgent"

    @classmethod
    def from_file(cls, file_path: str, config: Config):
        prs = PPTXPre(file_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        slides = []
        error_history = []
        slide_idx = 0
        layouts = [layout.name for layout in prs.slide_layouts]
        for slide in prs.slides:
            if slide._element.get("show") == "0":
                continue

            slide_idx += 1
            try:
                if slide.slide_layout.name not in layouts:
                    raise ValueError(
                        f"slide layout {slide.slide_layout.name} not found"
                    )
                slides.append(
                    SlidePage.from_slide(
                        slide,
                        slide_idx - len(error_history),
                        slide_idx,
                        slide_width.pt,
                        slide_height.pt,
                        config,
                    )
                )
            except Exception as e:
                error_history.append((slide_idx, str(e)))
                if config.DEBUG:
                    print(f"Warning in slide {slide_idx}: {e}")

        num_pages = len(slides)
        return cls(
            slides, error_history, slide_width, slide_height, file_path, num_pages
        )

    def save(self, file_path, layout_only=False):
        self.clear_slides()
        for slide in self.slides:
            if layout_only:
                self.clear_images(slide.shapes)
            pptx_slide = self.build_slide(slide)
            if layout_only:
                self.clear_text(pptx_slide.shapes)
        self.prs.save(file_path)

    def build_slide(self, slide: SlidePage) -> PPTXSlide:
        return slide.build(
            self.prs.slides.add_slide(self.layout_mapping[slide.slide_layout_name])
        )

    def clear_slides(self):
        while len(self.prs.slides) != 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]

    def clear_images(self, shapes: list[ShapeElement]):
        for idx, shape in enumerate(shapes):
            if isinstance(shape, GroupShape):
                self.clear_images(shape.data)
            elif isinstance(shape, Picture):
                shape.img_path = "resource/pic_placeholder.png"
            elif isinstance(shape, GraphicalShape):
                shapes[idx] = shape.normalize()

    def clear_text(self, shapes: list[BaseShape]):
        for shape in shapes:
            if isinstance(shape, PPTXGroupShape):
                self.clear_text(shape.shapes)
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = "a" * len(run.text)

    def to_html(self, **kwargs) -> str:
        return "\n".join(
            [
                f"Slide Page {slide_idx}\n" + slide.to_html(**kwargs)
                for slide_idx, slide in enumerate(self.slides)
            ]
        )

    def to_text(self, show_image: bool = True) -> str:
        return "\n".join(
            [
                f"Slide Page {slide_idx}\n" + slide.to_text(show_image)
                for slide_idx, slide in enumerate(self.slides)
            ]
        )

    def normalize(self):
        for slide in self.slides:
            slide.normalize()
        return self

    def __len__(self):
        return len(self.slides)


SHAPECAST: dict[int, ShapeElement] = {
    MSO_SHAPE_TYPE.AUTO_SHAPE: FreeShape,
    MSO_SHAPE_TYPE.PLACEHOLDER: Placeholder,
    MSO_SHAPE_TYPE.PICTURE: Picture,
    MSO_SHAPE_TYPE.GROUP: GroupShape,
    MSO_SHAPE_TYPE.TEXT_BOX: TextBox,
    MSO_SHAPE_TYPE.CHART: GraphicalShape,
    MSO_SHAPE_TYPE.TABLE: GraphicalShape,
    MSO_SHAPE_TYPE.DIAGRAM: GraphicalShape,
    MSO_SHAPE_TYPE.LINE: Connector,
}

if __name__ == "__main__":
    from glob import glob

    import tiktoken

    encoder = tiktoken.encoding_for_model("gpt-4o")

    for ppt in glob("data/*/pptx/*/source_standard.pptx"):
        prs = Presentation.from_file(ppt, Config("/tmp"))
        for slide in prs.slides:
            # geometry, size, area, font_style, element_id(image_id), textframe_id
            html = slide.to_html(
                geometry=True,
                size=True,
                area=True,
                element_id=True,
                textframe_id=True,
                font_style=True,
            )
            print("\033c", html, "total tokens:", len(encoder.encode(html)))
