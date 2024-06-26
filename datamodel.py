from dataclasses import dataclass
import json
from typing import List
from pptx import Presentation
from datetime import datetime


# meta class
class Page:
    def to_slide(self, ppt: Presentation) -> Presentation:
        return ppt


@dataclass
class Cover(Page):
    DepartmentName: str
    Topic: str
    Presenter: str

    def to_slide(self, ppt: Presentation):
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        placeholders = list(slide.shapes.placeholders)
        placeholders[0].text = self.Topic
        placeholders[1].text = self.DepartmentName
        placeholders[2].text = (
            self.Presenter + "\n" + datetime.now().strftime("%Y年%m月%d日")
        )
        return ppt


@dataclass
class TableOfContents(Page):
    Sections: List[str]

    def to_slide(self, ppt: Presentation):
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        placeholders = list(slide.shapes.placeholders)
        placeholders[0].text = "".join(
            [f" {sec}\n" for idx, sec in enumerate(self.Sections)]
        )


section_count = 0


@dataclass
class Section(Page):
    SectionName: str
    SubSections: List[str]

    def to_slide(self, ppt: Presentation):
        global section_count
        section_count += 1
        slide = ppt.slides.add_slide(ppt.slide_layouts[2])
        placeholders = list(slide.shapes.placeholders)
        placeholders[0].text = f"{section_count:02d}"
        placeholders[1].text = self.SectionName
        placeholders[2].text = "\n".join(self.SubSections)
        return ppt


@dataclass
class Highlight(Page):
    SubSectionName: str
    Highlights: List[str]
    Contents: str

    def to_slide(self, ppt: Presentation):
        slide = ppt.slides.add_slide(ppt.slide_layouts[3])
        placeholders = list(slide.shapes.placeholders)
        placeholders[4].text = self.Contents
        placeholders[0].text = self.SubSectionName
        if len(self.Highlights) >= 3:
            placeholders[2].text = self.Highlights[0]
            placeholders[3].text = self.Highlights[1]
            placeholders[1].text = self.Highlights[2]
        else:
            raise ("Error length of highlights")
        return ppt


@dataclass
class MultilevelBulletedList(Page):
    SubSectionName: str
    ContentsList: List[str]

    def to_slide(self, ppt: Presentation):
        slide = ppt.slides.add_slide(ppt.slide_layouts[4])
        placeholders = list(slide.shapes.placeholders)

        placeholders[0].text = self.SubSectionName
        for datapoint in self.ContentsList:
            placeholders[1].text += "● " + datapoint["topic"]+'\n' + "".join([f"     {i}\n" for i in datapoint["text"]])
        return ppt


@dataclass
class Display(Page):
    SubSectionName: str
    PicturePath: str
    PicTopic: str
    Contents: List[str]

    def to_slide(self, ppt: Presentation):
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])
        placeholders = list(slide.shapes.placeholders)
        placeholders[0].text = self.SubSectionName
        placeholders[2].text = "\n".join(self.Contents)

        placeholder = placeholders[1]
        placeholders[3].text = self.PicTopic
        placeholder.element.getparent().remove(placeholder.element)
        slide.shapes.add_picture(
            self.PicturePath,
            placeholder.left,
            placeholder.top,
            placeholder.width,
            placeholder.height,
        )
        return ppt


@dataclass
class Information(Page):
    SubSectionName: str
    InformationTopic: str
    Information: List[str]

    def to_slide(self, ppt: Presentation):
        slide = ppt.slides.add_slide(ppt.slide_layouts[6])
        placeholders = list(slide.shapes.placeholders)

        text = self.InformationTopic + "\n" + "\n".join(self.Information)
        placeholders[2].text = text
        placeholders[1].text = self.InformationTopic
        placeholders[0].text = self.SubSectionName
        return ppt


@dataclass
class End(Page):
    def to_slide(self, ppt: Presentation):
        ppt.slides.add_slide(ppt.slide_layouts[-1])
        return ppt

DATACLASSES= [ Cover, TableOfContents, Section, Highlight, MultilevelBulletedList, Display, Information, End ]

def json_to_dataclass(data: list) -> List[Page]:
    dataclass_list = []

    for page in data:
        # 获取全局命名空间的数据类
        DataClass = globals()[page["PageType"]]
        page.pop("PageType")
        dataclass_list.append(DataClass(**page))

    return dataclass_list
