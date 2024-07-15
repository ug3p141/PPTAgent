import pickle
import os
import re
from rich import print
from pptx.oxml import parse_xml
from pptx import Presentation as PPTXPre
from pptx.slide import Slide
from pptx.shapes.autoshape import Shape as PPTXAutoShape
from pptx.shapes.picture import Picture
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.chart.chart import Chart
from pptx.table import Table
from pptx.text.text import _Run, TextFrame
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.fill import _NoneFill, _NoFill
from typing_extensions import Dict, Type
from utils import (
    Config,
    dict_to_object,
    parse_groupshape,
    replace_xml_node,
    xml_print,
    object_to_dict,
)  # noqa
from pptx.dml.color import RGBColor
from pptx.util import Emu

base_config = Config()
pjoin = os.path.join


# 三种级别的text 可以保存的属性
# textframe: shape bounds
# paragraph: space, alignment, level, font
# run: font, hyperlink, text
# 以paragraphs 为单位处理文本框
class TextFrame:
    def __init__(
        self,
        is_textframe: bool,
        style: Dict,
        data: Dict,
    ):
        self.is_textframe = is_textframe
        self.style = style
        self.data = data
        self.descr = "".join([para.get("text", "") for para in data])

    @classmethod
    def from_shape(cls, shape: BaseShape):
        if not shape.has_text_frame:
            return cls(False, {}, [])
        shape = shape.text_frame
        # always set word wrap
        style = object_to_dict(shape, exclude=["text"])
        data = []
        # 获取文本框内容
        for _, paragraph in enumerate(shape.paragraphs):
            runs = paragraph.runs
            if len(runs) == 0:  # 成功解析
                runs = [
                    _Run(r, paragraph)
                    for r in parse_xml(paragraph._element.xml.replace("fld", "r")).r_lst
                ]
            # 只有fill属性懒得解析了
            content = object_to_dict(paragraph, exclude=["runs"])
            if len(runs) == 0:
                data.append(content)
                continue
            runs_ronts = {run.font: 0 for run in runs}
            for run in runs:
                runs_ronts[run.font] += len(run.text)
            content["font"] = object_to_dict(max(runs_ronts, key=runs_ronts.get))

            # 设置language id会导致font name设置失败，fuck！
            if content["font"]["name"] == "+mj-ea":
                content["font"]["name"] = "宋体"

            data.append(content)
        return cls(True, style, data)

    def build(self, shape: BaseShape):
        if not self.is_textframe:
            return
        tf = (
            shape.text_frame
            if shape.has_text_frame
            else shape.add_textbox(**self.style.pop("shape_bounds"))
        )
        # 注意，只对textframe是无效的，还要设置一遍所有的runs
        for pid, para in enumerate(self.data):
            text = para.pop("text")
            font = para.pop("font") if "font" in para else None

            p = tf.paragraphs[0]
            if pid != 0:
                p = tf.add_paragraph()
            dict_to_object(para, p)

            # 目前只有一个run，后续使用markdown的可能更多
            if font is None:
                continue
            run = p.add_run()
            if font["color"] is not None:
                run.font.fill.solid()
                run.font.fill.fore_color.rgb = RGBColor(*font.pop("color"))
            run.text = text
            dict_to_object(font, run.font)
        dict_to_object(self.style, tf)

    def __repr__(self) -> str:
        return f"{self.__class__.__name__}: {self.descr}"


# TODO 判断build后的属性是否与解析前的一致


class UnsupportedShape:
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: BaseShape,
        style: Dict,
        text_frame: TextFrame,
    ):
        # print(f"Unsupport Shape --- {shape.__class__}\n", object_to_dict(shape))
        assert not any([shape.has_chart, shape.has_table, shape.has_text_frame])
        obj = cls()
        obj.style = style
        return obj

    def build(self, slide: Slide):
        pass


class ShapeElement:
    def __init__(
        self,
        slide_idx: int,
        shape_idx: int,
        style: Dict,
        data: Dict,
        text_frame: TextFrame,
        descr: str = "",
    ):
        self.slide_idx = slide_idx
        self.shape_idx = shape_idx
        self.style = style
        self.data = data
        self.text_frame = text_frame
        self.validation = True
        self.consumed = False
        self.descr = descr

    @classmethod
    def from_shape(cls, slide_idx: int, shape_idx: int, shape: BaseShape):
        # implement fill for shapes

        style = {
            "shape_bounds": {
                "width": Emu(shape.width.emu),
                "height": Emu(shape.height.emu),
                "left": Emu(shape.left.emu),
                "top": Emu(shape.top.emu),
            },
            "shape_type": str(shape.shape_type),
            "rotation": shape.rotation,
            "fill": shape.fill._xPr.xml if "fill" in dir(shape) else None,
            "line": (
                {
                    "line_xml": shape.line.fill._xPr.xml,
                    "width": shape.line.width,
                    "dash_style": shape.line.dash_style,
                }
                if "line" in dir(shape) and shape.line._ln is not None
                else None
            ),
        }
        text_frame = TextFrame.from_shape(shape)
        obj = SHAPECAST.get(shape.shape_type, UnsupportedShape).from_shape(
            slide_idx, shape_idx, shape, style, text_frame
        )
        # ? mask to enable pickling
        obj.from_shape = shape
        if shape.is_placeholder:
            print("placeholder", obj)
        return obj

    def build(self, slide, shape):
        assert not self.consumed, "Shape has been consumed"
        self.consumed = True
        self.text_frame.build(shape)
        if self.style["fill"] is not None:
            replace_xml_node(shape.fill._xPr, self.style["fill"])
            dict_to_object(self.style["shape_bounds"], shape)
        if self.style["line"] is not None:
            # shape.line.fill.background()
            replace_xml_node(
                shape.line.fill._xPr,
                self.style["line"].pop("line_xml"),
            )
            dict_to_object(self.style["line"], shape.line)
        if "rotation" in dir(shape):
            shape.rotation = self.style["rotation"]
        return shape

    def get_style(self):
        return self.style

    def get_data(self):
        return self.data

    def __eq__(self, __value: object) -> bool:
        if not isinstance(__value, ShapeElement) or self.consumed:
            return False
        return (
            self.shape_idx == __value.shape_idx
            and self.style == __value.style
            and self.data == __value.data
            and self.text_frame == __value.text_frame
        )

    def __repr__(self) -> str:
        return f"{self.__class__.__name__}: shape {self.shape_idx} of slide {self.slide_idx}, {self.descr or self.text_frame.descr}"

    def pickle(self, path: str):
        # pickle self for later use
        with open(path, "wb") as f:
            pickle.dump(self, f)


class AutoShape(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: PPTXAutoShape,
        style: Dict,
        text_frame: TextFrame,
    ):
        data = {
            "auto_shape_type": shape.auto_shape_type.real,
            "is_nofill": isinstance(shape.fill._fill, (_NoneFill, _NoFill)),
            "is_line_nofill": isinstance(shape.line.fill._fill, (_NoneFill, _NoFill)),
        }
        descr = str(shape.auto_shape_type)
        return cls(slide_idx, shape_idx, style, data, text_frame, descr)

    def build(self, slide: Slide):
        shape = slide.shapes.add_shape(
            self.data["auto_shape_type"], **self.style["shape_bounds"]
        )
        if self.data["is_nofill"]:
            shape.fill.background()
            self.style["fill"] = None
            for para in self.text_frame.data:
                if "font" in para:
                    para["font"]["color"] = (0, 0, 0)
        if self.data["is_line_nofill"]:
            shape.line.fill.background()
            self.style["line"] = None
        shape = super().build(slide, shape)


class TextBox(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: TextFrame,
        style: Dict,
        text_frame: TextFrame,
    ):
        return cls(slide_idx, shape_idx, style, None, text_frame, False)

    def build(self, slide: Slide):
        shape = slide.shapes.add_textbox(**self.style["shape_bounds"])
        return super().build(slide, shape)


class Placeholder(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: SlidePlaceholder,
        style: Dict,
        text_frame: TextFrame,
    ):
        style |= {
            "placeholder_type": shape.placeholder_format.type,
        }
        if not shape.has_text_frame or shape.has_chart or shape.has_table:
            print(f"Unsupported Shape: {shape}")
        data = []
        return cls(slide_idx, shape_idx, style, data, text_frame)

    def build(self, slide: Slide):
        pass
        # super().build(slide, self.shape)


class Picture(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: Picture,
        style: Dict,
        text_frame: TextFrame,
    ):
        img_name = re.sub(r"[\/\0]", "_", shape.name)
        img_path = pjoin(
            base_config.IMAGE_DIR,
            f"{img_name}.{shape.image.ext}",
        )
        style["img_style"] = {
            "crop_bottom": shape.crop_bottom,
            "crop_top": shape.crop_top,
            "crop_left": shape.crop_left,
            "crop_right": shape.crop_right,
            "auto_shape_type": shape.auto_shape_type,  # most are rectangles
        }
        with open(img_path, "wb") as f:
            f.write(shape.image.blob)
        return cls(
            slide_idx, shape_idx, style, [img_path, shape.name], text_frame, shape.name
        )

    def build(self, slide: Slide):
        shape = slide.shapes.add_picture(
            self.data[0],
            **self.style["shape_bounds"],
        )
        shape.name = self.data[1]
        dict_to_object(self.style["img_style"], shape.image)
        return super().build(slide, shape)


class GroupShape(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: GroupShape,
        style: Dict,
        text_frame: TextFrame,
    ):
        data = [
            ShapeElement.from_shape(slide_idx, f"{shape_idx}_{i}", sub_shape)
            for i, sub_shape in enumerate(shape.shapes)
        ]
        for idx, shape_bounds in enumerate(parse_groupshape(shape)):
            data[idx].style["shape_bounds"] = shape_bounds
        return cls(slide_idx, shape_idx, style, data, text_frame, " group ")

    def build(self, slide: Slide):
        shapes = []
        for sub_shape in self.data:
            shapes.append(sub_shape.build(slide))
        return super().build(slide, shapes)


class Chart(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: Chart,
        style: Dict,
        text_frame: TextFrame,
    ):
        chart = shape.chart
        style["chart_style"] = object_to_dict(chart) | {
            "chart_title": (
                chart.chart_title.text_frame.text
                if chart.chart_title.has_text_frame
                else None
            ),
        }
        chart_data = []
        for series in chart.series:
            chart_item_data = {"name": str(series.name)}
            data_list = []
            for point in series.values:
                data_list.append(point)
            chart_item_data["data_list"] = data_list
            chart_data.append(chart_item_data)
        obj = cls(slide_idx, shape_idx, style, chart_data, text_frame)
        setattr(obj, "shape", shape)
        return obj

    def build(self, slide: Slide):
        pass
        # graphic_frame = slide.shapes.add_chart(
        #     self.style["chart_style"]["chart_type"],

        # shape
        # super().build(slide, self.shape)


# 目前的问题，rows和cols的宽度或高度获取不到
class Table(ShapeElement):
    @classmethod
    def from_shape(
        cls,
        slide_idx: int,
        shape_idx: int,
        shape: Table,
        style: Dict,
        text_frame: TextFrame,
    ):
        table = shape.table
        row_list = list()
        for row in table.rows:
            cell_list = list()
            for cell in row.cells:
                cell_data = cell.text_frame.text
                cell_list.append(cell_data)
            row_list.append(cell_list)
        style["table_style"] = object_to_dict(shape.table)

        obj = cls(slide_idx, shape_idx, style, row_list, text_frame)
        setattr(obj, "shape", shape)
        return obj

    def build(self, slide: Slide):
        graphic_frame = slide.shapes.add_table(
            len(self.data), len(self.data[0]), **self.style["shape_bounds"]
        )

        dict_to_object(self.style["table_style"], graphic_frame.table)
        self.text_frame.build(graphic_frame)


# template包含哪些东西，如何定义一个template
# 1. template pics：例如banner、background pic、logo，此部分只需要手动识别插入以及build
# 2. cloneable shapes：例如listed text, parallel text等，此部分需要手动安排布局
class SlidePage:
    def __init__(
        self,
        shapes: list[ShapeElement],
        slide_idx: int,
        slide_notes: str,
        slide_layout_name: str,
        slide_title: str,
    ):
        self.slide_idx = slide_idx
        self.shapes = shapes
        self.slide_layout_name = slide_layout_name
        self.slide_title = slide_title
        self.consumed = False

    @classmethod
    def from_slide(cls, slide: Slide, slide_idx: int):
        shapes = [
            ShapeElement.from_shape(slide_idx, i, shape)
            for i, shape in enumerate(slide.shapes)
        ]
        #! cannot be sort because it will change the z-order
        # shapes.sort(key=lambda x: x.style["shape_bounds"]["top"])
        slide_layout_name = slide.slide_layout.name if slide.slide_layout else None
        # is shape a title
        slide_title = slide.shapes.title.text if slide.shapes.title else None
        slide_notes = (
            slide.notes_slide.notes_text_frame.text
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame
            else None
        )
        return cls(shapes, slide_idx, slide_notes, slide_layout_name, slide_title)

    def build(self, prs: PPTXPre, slide_layout):
        assert not self.consumed, "SlidePage has been consumed"
        self.consumed = True
        slide = prs.slides.add_slide(slide_layout)
        for shape in self.shapes:
            shape.build(slide)
        # if self.slide_title:
        # slide.shapes.title.text = self.slide_title

    def __eq__(self, __value) -> bool:
        if not isinstance(__value, SlidePage) or self.consumed:
            return False
        return (
            self.slide_idx == __value.slide_idx
            and self.shapes == __value.shapes
            and self.slide_layout_name == __value.slide_layout_name
            and self.slide_title == __value.slide_title
        )


class Presentation:
    def __init__(
        self,
        slides: list[SlidePage],
        slide_width: float,
        slide_height: float,
        file_path: str,
        num_pages: int,
        author: str,
    ) -> None:
        self.slides = slides
        self.slide_width = slide_width
        self.slide_height = slide_height
        self.num_pages = num_pages
        self.source_file = file_path
        self.author = author
        self.consumed = False

    @classmethod
    def from_file(cls, file_path):
        prs = PPTXPre(file_path)
        # TODO set slide range here
        slides = [
            SlidePage.from_slide(slide, i)
            for i, slide in enumerate(list(prs.slides)[:])
        ]
        slide_width = Emu(prs.slide_width.emu)
        slide_height = Emu(prs.slide_height.emu)
        core_properties = prs.core_properties
        num_pages = len(slides)
        author = core_properties.author
        return cls(slides, slide_width, slide_height, file_path, num_pages, author)

    def save(self, file_path, validation=False):
        assert not self.consumed, "Presentation has been consumed"
        if validation:
            copied_pr = Presentation.from_file(self.file_path)
        self.consumed = True
        prs = PPTXPre(self.source_file)
        # delete all slides so we can rebuild them
        while len(prs.slides) != 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        prs.core_properties.author += "OminiPreGen with: "
        layout_mapping = {layout.name: layout for layout in prs.slide_layouts}
        for slide in self.slides:
            slide.build(prs, layout_mapping[slide.slide_layout_name])
        if validation:
            saved_prs = Presentation.from_file(file_path)
            assert copied_pr == saved_prs, "Validation failed"
        prs.save(file_path)

    def __eq__(self, __value: object) -> bool:
        if not isinstance(__value, Presentation) or self.consumed:
            return False
        slides_eq = True
        for i, slide in enumerate(self.slides):
            if not slide == __value.slides[i]:
                slides_eq = False
                break
        return (
            slides_eq
            and self.slide_width == __value.slide_width
            and self.slide_height == __value.slide_height
            and self.num_pages == __value.num_pages
        )


SHAPECAST: Dict[int, Type[ShapeElement]] = {
    MSO_SHAPE_TYPE.AUTO_SHAPE: AutoShape,
    MSO_SHAPE_TYPE.PLACEHOLDER: Placeholder,
    MSO_SHAPE_TYPE.PICTURE: Picture,
    MSO_SHAPE_TYPE.GROUP: GroupShape,
    MSO_SHAPE_TYPE.CHART: Chart,
    MSO_SHAPE_TYPE.TABLE: Table,
    MSO_SHAPE_TYPE.TEXT_BOX: TextBox,
}


if __name__ == "__main__":
    Presentation.from_file(
        pjoin(
            base_config.PPT_DIR,
            "中文信息联合党支部2022年述职报告.pptx",
        )
    ).save(pjoin(base_config.GEN_PPT_DIR, "ppt_handlers_test.pptx"))
